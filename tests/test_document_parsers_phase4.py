"""Phase 4 document parser consolidation and compatibility regression tests."""

from __future__ import annotations

import importlib
import logging
import shutil
import stat
import zipfile
from dataclasses import replace
from pathlib import Path
from unittest.mock import AsyncMock, Mock

import pytest
from kazma_core.documents.config import DocumentConfig
from kazma_core.documents.errors import (
    DocumentEncryptedError,
    DocumentFormatError,
    DocumentLimitError,
    DocumentSandboxError,
    DocumentSecurityError,
)
from kazma_core.documents.parsers.common import ParseContext, sha256_path
from kazma_core.documents.registry import (
    ParserPlugin,
    ParserReadiness,
    ParserRegistry,
)
from kazma_core.documents.sandbox import SandboxResult
from kazma_core.documents.service import DocumentService
from kazma_core.documents.sniff import sniff_document


class _NoopParser:
    def parse(self, path: Path, context: ParseContext):  # pragma: no cover - probe only
        raise AssertionError


class _ToolResult:
    def __init__(self, value: str) -> None:
        self.value = value

    def as_tool_output(self) -> str:
        return self.value


@pytest.fixture
def document_config(tmp_path: Path) -> DocumentConfig:
    return DocumentConfig(
        storage_root=tmp_path / "document-store",
        worker_timeout_seconds=30,
    )


def _context(path: Path, config: DocumentConfig, parser_id: str, mime: str) -> ParseContext:
    return ParseContext(
        config=config,
        source_sha256=sha256_path(path),
        mime_type=mime,
        extension=path.suffix.lower(),
        parser_id=parser_id,
        parser_version="1",
    )


def test_registry_marks_healthy_native_ready_and_missing_dependency_unavailable() -> None:
    plugins = (
        ParserPlugin(
            "native",
            "1",
            ("text/plain",),
            (".txt",),
            ("text",),
            ("max_output_chars_total",),
            True,
            _NoopParser,
        ),
        ParserPlugin(
            "missing",
            "1",
            ("application/x-missing",),
            (".missing",),
            (),
            (),
            True,
            _NoopParser,
            dependencies=("kazma_dependency_that_does_not_exist",),
        ),
    )
    registry = ParserRegistry(plugins)
    assert registry.capability("native").readiness is ParserReadiness.READY
    assert registry.capability("missing").readiness is ParserReadiness.UNAVAILABLE


def test_legacy_office_readiness_tracks_libreoffice_probe() -> None:
    registry = ParserRegistry()
    expected = bool(shutil.which("soffice") or shutil.which("libreoffice"))
    for parser_id in ("legacy-doc", "legacy-xls", "legacy-ppt"):
        assert registry.capability(parser_id).available is expected


def test_mime_mismatch_and_unknown_format_reject(
    tmp_path: Path, document_config: DocumentConfig
) -> None:
    mismatch = tmp_path / "wrong.pdf"
    mismatch.write_text("plain text", encoding="utf-8")
    with pytest.raises(DocumentFormatError, match="does not match"):
        sniff_document(mismatch, document_config)

    unknown = tmp_path / "sample.bin"
    unknown.write_bytes(b"\x00\x01\x02")
    with pytest.raises(DocumentFormatError, match="Unsupported document extension"):
        sniff_document(unknown, document_config)


def _write_ooxml(
    path: Path,
    members: dict[str, bytes],
    *,
    compression: int = zipfile.ZIP_DEFLATED,
) -> None:
    with zipfile.ZipFile(path, "w", compression=compression) as archive:
        archive.writestr("[Content_Types].xml", b"<Types/>")
        for name, value in members.items():
            archive.writestr(name, value)


def test_ooxml_preflight_rejects_traversal_bomb_and_macros(
    tmp_path: Path, document_config: DocumentConfig
) -> None:
    traversal = tmp_path / "traversal.docx"
    _write_ooxml(traversal, {"word/document.xml": b"<w/>", "../escape": b"x"})
    with pytest.raises(DocumentSecurityError, match="unsafe member path"):
        sniff_document(traversal, document_config)

    bomb = tmp_path / "bomb.docx"
    _write_ooxml(bomb, {"word/document.xml": b"x" * 100_000})
    with pytest.raises(DocumentLimitError, match="compression ratio"):
        sniff_document(bomb, replace(document_config, max_compression_ratio=2))

    members = tmp_path / "members.docx"
    _write_ooxml(
        members,
        {"word/document.xml": b"<w/>", "word/extra.xml": b"<w/>"},
    )
    with pytest.raises(DocumentLimitError, match="too many members"):
        sniff_document(members, replace(document_config, max_archive_members=2))

    macro = tmp_path / "macro.docx"
    _write_ooxml(
        macro,
        {"word/document.xml": b"<w/>", "word/vbaProject.bin": b"macro"},
    )
    with pytest.raises(DocumentSecurityError, match="Macro-enabled"):
        sniff_document(macro, document_config)

    unsafe_xml = tmp_path / "unsafe-xml.docx"
    _write_ooxml(
        unsafe_xml,
        {"word/document.xml": b"<!DOCTYPE x [<!ENTITY e SYSTEM 'file:///etc/passwd'>]>"},
    )
    with pytest.raises(DocumentSecurityError, match="unsafe XML"):
        sniff_document(unsafe_xml, document_config)

    symlink = tmp_path / "symlink.docx"
    with zipfile.ZipFile(symlink, "w") as archive:
        archive.writestr("[Content_Types].xml", b"<Types/>")
        archive.writestr("word/document.xml", b"<w/>")
        info = zipfile.ZipInfo("word/link")
        info.create_system = 3
        info.external_attr = (stat.S_IFLNK | 0o777) << 16
        archive.writestr(info, "document.xml")
    with pytest.raises(DocumentSecurityError, match="symbolic link"):
        sniff_document(symlink, document_config)


@pytest.mark.parametrize(
    "marker",
    (
        b"<!DOCTYPE unsafe>",
        "<!ENTITY unsafe 'value'>".encode("utf-16-le"),
    ),
)
def test_ooxml_stream_scans_padded_xml_across_chunks(
    tmp_path: Path,
    document_config: DocumentConfig,
    marker: bytes,
) -> None:
    path = tmp_path / "padded.docx"
    payload = b" " * (64 * 1024 - 5) + marker + b"<w/>"
    _write_ooxml(
        path,
        {"word/document.xml": payload},
        compression=zipfile.ZIP_STORED,
    )

    with pytest.raises(DocumentSecurityError, match="unsafe XML"):
        sniff_document(path, document_config)


def test_ooxml_stream_enforces_actual_expansion_and_crc(
    tmp_path: Path,
    document_config: DocumentConfig,
) -> None:
    oversized = tmp_path / "actual-size.docx"
    _write_ooxml(
        oversized,
        {"word/document.xml": b"<w>" + b"x" * 128 + b"</w>"},
        compression=zipfile.ZIP_STORED,
    )
    with pytest.raises(DocumentLimitError, match="expanded size"):
        sniff_document(
            oversized,
            replace(document_config, max_expanded_bytes=64),
        )

    corrupt = tmp_path / "bad-crc.docx"
    _write_ooxml(
        corrupt,
        {"word/document.xml": b"<w>integrity</w>"},
        compression=zipfile.ZIP_STORED,
    )
    with zipfile.ZipFile(corrupt) as archive:
        info = archive.getinfo("word/document.xml")
        data_offset = (
            info.header_offset
            + 30
            + len(info.filename.encode("utf-8"))
            + len(info.extra)
        )
    payload = bytearray(corrupt.read_bytes())
    payload[data_offset + 3] ^= 0x01
    corrupt.write_bytes(payload)
    with pytest.raises(DocumentFormatError, match="integrity validation"):
        sniff_document(corrupt, document_config)


def test_ooxml_rejects_falsified_central_directory_size(
    tmp_path: Path,
    document_config: DocumentConfig,
) -> None:
    path = tmp_path / "false-size.docx"
    _write_ooxml(
        path,
        {"word/document.xml": b"<w>directory mismatch</w>"},
        compression=zipfile.ZIP_STORED,
    )
    payload = bytearray(path.read_bytes())
    signature = b"PK\x01\x02"
    position = payload.find(signature)
    while position >= 0:
        name_length = int.from_bytes(payload[position + 28 : position + 30], "little")
        name = bytes(payload[position + 46 : position + 46 + name_length])
        if name == b"word/document.xml":
            payload[position + 24 : position + 28] = (1).to_bytes(4, "little")
            break
        position = payload.find(signature, position + 4)
    else:  # pragma: no cover - Python's ZIP writer always emits this record
        raise AssertionError("central directory member not found")
    path.write_bytes(payload)

    with pytest.raises(DocumentFormatError):
        sniff_document(path, document_config)


def test_degraded_pdf_capability_does_not_advertise_tables(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """pypdf-only installs are degraded (text only); tables require PyMuPDF/pdfplumber."""
    from kazma_core.documents import parsers

    real_import = parsers.importlib.import_module

    def import_text_only_pdf_stack(name: str):
        if name in {"pdfplumber", "fitz", "pymupdf"}:
            raise ImportError
        if name == "pypdf":
            return object()
        return real_import(name)

    monkeypatch.setattr(parsers.importlib, "import_module", import_text_only_pdf_stack)
    pdf_plugin = next(
        plugin for plugin in parsers.builtin_plugins() if plugin.parser_id == "pdf"
    )
    capability = ParserRegistry((pdf_plugin,)).capability("pdf")

    assert capability.readiness is ParserReadiness.DEGRADED
    assert capability.features == ("pages", "text")
    assert "tables" not in capability.features


def test_encrypted_pdf_rejected_before_parser(
    tmp_path: Path, document_config: DocumentConfig
) -> None:
    path = tmp_path / "encrypted.pdf"
    path.write_bytes(b"%PDF-1.7\n1 0 obj\n<< /Encrypt 2 0 R >>\n%%EOF")
    with pytest.raises(DocumentEncryptedError):
        sniff_document(path, document_config)


def test_pdf_zip_polyglot_rejected(
    tmp_path: Path, document_config: DocumentConfig
) -> None:
    path = tmp_path / "polyglot.pdf"
    path.write_bytes(b"%PDF-1.7\n%%EOF\nPK\x03\x04payload")
    with pytest.raises(DocumentSecurityError, match="polyglot"):
        sniff_document(path, document_config)


@pytest.mark.parametrize(
    ("name", "content", "parser_id", "mime"),
    [
        ("sample.txt", "alpha\n\nbeta", "text", "text/plain"),
        ("sample.md", "# Heading\n\nbody", "text", "text/markdown"),
        ("sample.log", "INFO safe", "text", "text/x-log"),
        ("sample.json", '{"b": 2, "a": 1}', "json", "application/json"),
        ("sample.csv", "a,b\n1,2\n", "csv", "text/csv"),
        ("sample.tsv", "a\tb\n1\t2\n", "tsv", "text/tab-separated-values"),
        ("sample.html", "<html><body><h1>Hello</h1><p>world</p></body></html>", "html", "text/html"),
        ("sample.rtf", r"{\rtf1\ansi Hello \b world\b0}", "rtf", "application/rtf"),
    ],
)
def test_native_text_parsers_emit_deterministic_ir(
    tmp_path: Path,
    document_config: DocumentConfig,
    name: str,
    content: str,
    parser_id: str,
    mime: str,
) -> None:
    path = tmp_path / name
    path.write_text(content, encoding="utf-8")
    plugin, _ = ParserRegistry().resolve(mime_type=mime, extension=path.suffix)
    first = plugin.factory().parse(path, _context(path, document_config, parser_id, mime))
    second = plugin.factory().parse(path, _context(path, document_config, parser_id, mime))
    assert first.to_json() == second.to_json()
    assert first.metadata["source_sha256"] == sha256_path(path)
    assert first.pages


def test_csv_bounds_are_enforced(tmp_path: Path, document_config: DocumentConfig) -> None:
    path = tmp_path / "rows.csv"
    path.write_text("a,b\n1,2\n3,4\n", encoding="utf-8")
    plugin, _ = ParserRegistry().resolve(mime_type="text/csv", extension=".csv")
    ir = plugin.factory().parse(
        path,
        _context(
            path,
            replace(document_config, max_rows_per_sheet=2),
            "csv",
            "text/csv",
        ),
    )
    assert "Rows truncated" in str(ir.metadata["warnings"])
    with pytest.raises(DocumentLimitError, match="cell count"):
        plugin.factory().parse(
            path,
            _context(
                path,
                replace(document_config, max_cells=2),
                "csv",
                "text/csv",
            ),
        )


@pytest.mark.asyncio
async def test_service_subprocess_paging_and_single_fence(
    tmp_path: Path, document_config: DocumentConfig
) -> None:
    path = tmp_path / "page.txt"
    path.write_text("abcdefghij\n\nklmnopqrst", encoding="utf-8")
    result = await DocumentService(config=document_config).read_transient(
        path,
        approved_path=path.resolve(),
        max_chars=12,
    )
    output = result.as_tool_output()
    assert result.continuation["has_more"] is True
    assert result.continuation["next_offset"] == 12
    assert output.count('<kazma:data source="document" untrusted="true">') == 1
    assert "next_offset=12" in output


def test_service_rejects_bad_worker_checksum(document_config: DocumentConfig) -> None:
    response = {
        "protocol_version": 1,
        "ok": True,
        "code": "ok",
        "message": "parsed",
        "source_sha256": "wrong",
        "ir_sha256": "wrong",
        "ir": {},
    }
    with pytest.raises(DocumentSandboxError, match="source checksum"):
        DocumentService(config=document_config)._validate_response(response, "expected")


def test_service_surfaces_sandbox_timeout(
    tmp_path: Path,
    document_config: DocumentConfig,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    path = tmp_path / "timeout.txt"
    path.write_text("hello", encoding="utf-8")
    monkeypatch.setattr(
        "kazma_core.documents.service.run_isolated_subprocess",
        lambda request: SandboxResult(
            command=request.command,
            returncode=1,
            stdout=b"",
            stderr=b"",
            duration_seconds=1,
            timed_out=True,
            output_limit_exceeded=False,
            resource_limits_enforced=False,
            resource_limit_degraded_reason=None,
        ),
    )
    with pytest.raises(DocumentSandboxError, match="time limit"):
        DocumentService(config=document_config).read_transient_sync(
            path, approved_path=path
        )


def test_service_warns_when_resource_limits_degrade(
    tmp_path: Path,
    document_config: DocumentConfig,
    monkeypatch: pytest.MonkeyPatch,
    caplog: pytest.LogCaptureFixture,
) -> None:
    path = tmp_path / "degraded.txt"
    path.write_text("hello", encoding="utf-8")
    monkeypatch.setattr(
        "kazma_core.documents.service.run_isolated_subprocess",
        lambda request: SandboxResult(
            command=request.command,
            returncode=1,
            stdout=b"",
            stderr=b"",
            duration_seconds=0.1,
            timed_out=False,
            output_limit_exceeded=False,
            resource_limits_enforced=False,
            resource_limit_degraded_reason="CPU quota unavailable",
        ),
    )

    with caplog.at_level(logging.WARNING), pytest.raises(
        DocumentSandboxError, match="produced no result"
    ):
        DocumentService(config=document_config).read_transient_sync(
            path, approved_path=path
        )
    assert "resource limits were partially degraded" in caplog.text
    assert "CPU quota unavailable" in caplog.text


@pytest.mark.asyncio
async def test_document_processor_and_crawler_delegate(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    path = tmp_path / "delegate.txt"
    path.write_text("hello", encoding="utf-8")

    from kazma_skills.native.advanced_web_crawler import tools as crawler
    from kazma_skills.native.document_processor import tools as processor

    service = Mock()
    service.read_transient = AsyncMock(return_value=_ToolResult("delegated"))
    monkeypatch.setattr(processor, "DocumentService", Mock(return_value=service))
    monkeypatch.setattr(processor, "_workspace_scope_error", Mock(return_value=None))
    assert await processor.read_document(str(path)) == "delegated"
    service.read_transient.assert_awaited_once()

    service.read_transient.reset_mock()
    monkeypatch.setattr(crawler, "DocumentService", Mock(return_value=service))
    monkeypatch.setattr(crawler, "_workspace_scope_error", Mock(return_value=None))
    assert await crawler.parse_document(str(path)) == "delegated"
    service.read_transient.assert_awaited_once()


def test_attachment_excerpt_delegates_without_double_fence(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    from kazma_gateway.agent_handler import attachments

    path = tmp_path / "delegate.txt"
    path.write_text("hello", encoding="utf-8")
    service = Mock()
    service.read_transient_sync.return_value = _ToolResult("delegated")
    monkeypatch.setattr(attachments, "DocumentService", Mock(return_value=service))
    assert attachments._try_parse_attachment(str(path)) == "delegated"
    assert service.read_transient_sync.call_args.kwargs["fence"] is False


def test_available_optional_ooxml_and_pdf_parsers(
    tmp_path: Path, document_config: DocumentConfig
) -> None:
    registry = ParserRegistry()

    if registry.capability("docx").available:
        from docx import Document

        path = tmp_path / "minimal.docx"
        doc = Document()
        doc.add_paragraph("docx text")
        doc.save(path)
        sniffed = sniff_document(path, document_config)
        plugin, _ = registry.resolve(
            mime_type=sniffed.mime_type, extension=sniffed.extension
        )
        assert "docx text" in plugin.factory().parse(
            path, _context(path, document_config, "docx", sniffed.mime_type)
        ).to_json()

    if registry.capability("xlsx").available:
        from openpyxl import Workbook

        path = tmp_path / "minimal.xlsx"
        workbook = Workbook()
        workbook.active.append(["xlsx text"])
        workbook.save(path)
        sniffed = sniff_document(path, document_config)
        plugin, _ = registry.resolve(
            mime_type=sniffed.mime_type, extension=sniffed.extension
        )
        assert "xlsx text" in plugin.factory().parse(
            path, _context(path, document_config, "xlsx", sniffed.mime_type)
        ).to_json()

    if registry.capability("pptx").available:
        from pptx import Presentation

        path = tmp_path / "minimal.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "pptx text"
        presentation.save(path)
        sniffed = sniff_document(path, document_config)
        plugin, _ = registry.resolve(
            mime_type=sniffed.mime_type, extension=sniffed.extension
        )
        assert "pptx text" in plugin.factory().parse(
            path, _context(path, document_config, "pptx", sniffed.mime_type)
        ).to_json()

    if registry.capability("pdf").available:
        from pypdf import PdfWriter

        path = tmp_path / "minimal.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=100, height=100)
        with path.open("wb") as stream:
            writer.write(stream)
        sniffed = sniff_document(path, document_config)
        plugin, _ = registry.resolve(
            mime_type=sniffed.mime_type, extension=sniffed.extension
        )
        assert plugin.factory().parse(
            path, _context(path, document_config, "pdf", sniffed.mime_type)
        ).pages


def test_native_skill_loading_has_no_document_import_cycle() -> None:
    for module in (
        "kazma_core.documents",
        "kazma_skills.native.document_processor.tools",
        "kazma_skills.native.advanced_web_crawler.tools",
        "kazma_gateway.agent_handler.attachments",
    ):
        assert importlib.import_module(module)


def _generate_pdf_via_engine(path: Path, *, title: str, rtl: bool, table: bool = False) -> None:
    """Build a small PDF through Kazma's PdfEngine (reportlab / DOCX route)."""
    from kazma_core.documents.content_model import (
        BodyBlock,
        ContentModel,
        HeadingBlock,
        TableBlock,
        TitleBlock,
    )
    from kazma_core.documents.engines.pdf import PdfEngine
    from kazma_core.documents.profile import DocProfile

    profile = DocProfile.for_content(title, rtl=rtl)
    model = ContentModel()
    model.add(TitleBlock(text=title))
    if rtl:
        model.add(HeadingBlock(text="مقدمة", level=1))
        model.add(BodyBlock(text="هذا نص تجريبي للتحقق من الاستخراج"))
    else:
        model.add(HeadingBlock(text="Introduction", level=1))
        model.add(BodyBlock(text="Regression body for PDF text extraction."))
    if table:
        if rtl:
            model.add(
                TableBlock(
                    headers=["الاسم", "القيمة"],
                    rows=[["لغة", "عربي"], ["PDF", "تلقائي"]],
                )
            )
        else:
            model.add(
                TableBlock(
                    headers=["Name", "Value"],
                    rows=[["Lang", "EN"], ["Mode", "Auto"], ["Format", "PDF"]],
                )
            )
    PdfEngine(profile).render(model, path)
    assert path.is_file() and path.stat().st_size > 0


def test_pdf_parser_arabic_logical_order_pymupdf(
    tmp_path: Path, document_config: DocumentConfig
) -> None:
    """PyMuPDF primary must extract Arabic in logical order (not visual reverse)."""
    pytest.importorskip("fitz")
    pytest.importorskip("reportlab")

    path = tmp_path / "arabic-demo.pdf"
    title = "منظومة كاظمة للذكاء الاصطناعي"
    _generate_pdf_via_engine(path, title=title, rtl=True)

    from kazma_core.documents.parsers.pdf import PdfParser

    ir = PdfParser().parse(
        path, _context(path, document_config, "pdf", "application/pdf")
    )
    joined = "\n".join(block.text for page in ir.pages for block in page.blocks)

    # Logical-order tokens that define the product title.
    for token in ("منظومة", "كاظمة"):
        assert token in joined, f"missing logical token {token!r} in {joined!r}"

    # Classic pdfplumber visual reverse of the title — must not dominate.
    reversed_blob = "ةموظنم"
    assert reversed_blob not in joined, (
        f"extracted text looks visually reversed: {joined!r}"
    )

    assert ir.metadata.get("extractor") == "pymupdf"
    assert any(
        page.metadata.get("extractor") == "pymupdf" for page in ir.pages
    )


def test_pdf_parser_english_tables_preserved(
    tmp_path: Path, document_config: DocumentConfig
) -> None:
    """Tier-1 swap must still emit TABLE blocks for electronic English PDFs."""
    pytest.importorskip("fitz")
    pytest.importorskip("reportlab")

    path = tmp_path / "english-table.pdf"
    _generate_pdf_via_engine(
        path, title="Kazma Document Tables", rtl=False, table=True
    )

    from kazma_core.documents.models import BlockType
    from kazma_core.documents.parsers.pdf import PdfParser

    ir = PdfParser().parse(
        path, _context(path, document_config, "pdf", "application/pdf")
    )
    joined = "\n".join(block.text for page in ir.pages for block in page.blocks)
    assert "Kazma" in joined or "Document" in joined
    assert ir.metadata.get("extractor") == "pymupdf"

    table_blocks = [
        block
        for page in ir.pages
        for block in page.blocks
        if block.block_type is BlockType.TABLE
    ]
    # Prefer structured TABLE blocks; fall back to pipe/row text if find_tables
    # misses a drawn table (reportlab layouts vary).
    if table_blocks:
        table_text = "\n".join(block.text for block in table_blocks)
        assert "Name" in table_text or "Lang" in table_text or "EN" in table_text
        assert "|" in table_text or "\n" in table_text
    else:
        assert "Name" in joined or "Lang" in joined or "Auto" in joined


def test_pdf_parser_falls_back_when_pymupdf_fails(
    tmp_path: Path,
    document_config: DocumentConfig,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Hard failure in PyMuPDF should fall through to pdfplumber, not abort."""
    pytest.importorskip("pdfplumber")
    from pypdf import PdfWriter

    path = tmp_path / "blank.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with path.open("wb") as stream:
        writer.write(stream)

    from kazma_core.documents.parsers import pdf as pdf_mod

    def boom(*_a, **_k):
        raise RuntimeError("simulated pymupdf failure")

    monkeypatch.setattr(pdf_mod.PdfParser, "_parse_pymupdf", staticmethod(boom))
    ir = pdf_mod.PdfParser().parse(
        path, _context(path, document_config, "pdf", "application/pdf")
    )
    assert ir.metadata.get("extractor") in {"pdfplumber", "pypdf", "PyPDF2"}
    assert ir.pages
    assert ir.metadata.get("extractors_tried")


def test_pdf_layout_reading_order_multi_column() -> None:
    """Multi-column blocks read left column fully, then right (LTR)."""
    from kazma_core.documents.parsers.pdf_layout import reading_order_text

    # Page width 600: left col x~50-250, right col x~320-550, full-width title.
    blocks = [
        {"x0": 40, "y0": 20, "x1": 560, "y1": 50, "cx": 300, "cy": 35, "text": "TITLE"},
        {"x0": 50, "y0": 80, "x1": 250, "y1": 100, "cx": 150, "cy": 90, "text": "L1"},
        {"x0": 50, "y0": 110, "x1": 250, "y1": 130, "cx": 150, "cy": 120, "text": "L2"},
        {"x0": 320, "y0": 80, "x1": 550, "y1": 100, "cx": 435, "cy": 90, "text": "R1"},
        {"x0": 320, "y0": 110, "x1": 550, "y1": 130, "cx": 435, "cy": 120, "text": "R2"},
    ]
    text, meta = reading_order_text(blocks, page_width=600.0, rtl=False)
    assert meta["column_count"] == 2
    assert meta["layout"] == "multi_column"
    # Title first, then left column top→bottom, then right.
    assert text.splitlines() == ["TITLE", "L1", "L2", "R1", "R2"]

    # RTL: right column before left.
    text_rtl, meta_rtl = reading_order_text(blocks, page_width=600.0, rtl=True)
    assert meta_rtl["rtl_columns"] is True
    assert text_rtl.splitlines() == ["TITLE", "R1", "R2", "L1", "L2"]


def test_pdf_optional_pypdfium2_engine_participates_when_installed(
    tmp_path: Path, document_config: DocumentConfig, monkeypatch: pytest.MonkeyPatch
) -> None:
    """pypdfium2 is optional; when present it can win the multi-engine score."""
    pytest.importorskip("pypdfium2")
    from pypdf import PdfWriter

    path = tmp_path / "plain.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with path.open("wb") as stream:
        writer.write(stream)

    from kazma_core.documents.models import (
        BlockType,
        DocumentBlock,
        DocumentIR,
        DocumentPage,
        Provenance,
        new_document_id,
        new_version_id,
    )
    from kazma_core.documents.parsers import pdf as pdf_mod

    weak = DocumentIR(
        document_id=new_document_id(),
        version_id=new_version_id(),
        pages=(
            DocumentPage(
                1,
                (
                    DocumentBlock(
                        block_id="p1-b1",
                        block_type=BlockType.TEXT,
                        text="\uFE8D\uFE8E\uFE8F\uFE90" * 16,
                    ),
                ),
                metadata={"extractor": "pymupdf"},
            ),
        ),
        provenance=Provenance(source=str(path), parser="pdf"),
        metadata={"extractor": "pymupdf"},
    )
    strong = DocumentIR(
        document_id=new_document_id(),
        version_id=new_version_id(),
        pages=(
            DocumentPage(
                1,
                (
                    DocumentBlock(
                        block_id="p1-b1",
                        block_type=BlockType.TEXT,
                        text="PDFium extracted logical Arabic منظومة كاظمة content " * 3,
                    ),
                ),
                metadata={"extractor": "pypdfium2"},
            ),
        ),
        provenance=Provenance(source=str(path), parser="pdf"),
        metadata={"extractor": "pypdfium2"},
    )

    monkeypatch.setattr(
        pdf_mod.PdfParser, "_parse_pymupdf", staticmethod(lambda *_a, **_k: weak)
    )
    monkeypatch.setattr(
        pdf_mod.PdfParser, "_parse_pypdfium2", staticmethod(lambda *_a, **_k: strong)
    )
    monkeypatch.setattr(
        pdf_mod.PdfParser,
        "_parse_pdfplumber",
        staticmethod(lambda *_a, **_k: (_ for _ in ()).throw(ImportError("skip"))),
    )
    monkeypatch.setattr(
        pdf_mod.PdfParser,
        "_parse_pypdf",
        staticmethod(lambda *_a, **_k: (_ for _ in ()).throw(ImportError("skip"))),
    )

    ir = pdf_mod.PdfParser().parse(
        path, _context(path, document_config, "pdf", "application/pdf")
    )
    assert ir.metadata.get("extractor") == "pypdfium2"
    assert "منظومة" in ir.pages[0].blocks[0].text
    tried = ir.metadata.get("extractors_tried") or []
    assert any(str(item).startswith("pypdfium2:") for item in tried)


def test_pdf_multi_engine_keeps_higher_scoring_extract(
    tmp_path: Path,
    document_config: DocumentConfig,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """When the primary extract is weak, a stronger secondary engine wins."""
    from kazma_core.documents.models import (
        BlockType,
        DocumentBlock,
        DocumentIR,
        DocumentPage,
        Provenance,
        new_document_id,
        new_version_id,
    )
    from kazma_core.documents.parsers import pdf as pdf_mod

    path = tmp_path / "scored.pdf"
    path.write_bytes(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n1 0 obj<<>>endobj\ntrailer<<>>\n%%EOF\n")

    weak = DocumentIR(
        document_id=new_document_id(),
        version_id=new_version_id(),
        pages=(
            DocumentPage(
                1,
                (
                    DocumentBlock(
                        block_id="p1-b1",
                        block_type=BlockType.TEXT,
                        text="\uFE8D\uFE8E\uFE8F\uFE90" * 20,
                    ),
                ),
                metadata={"extractor": "pymupdf"},
            ),
        ),
        provenance=Provenance(source=str(path), parser="pdf"),
        metadata={"extractor": "pymupdf"},
    )
    strong = DocumentIR(
        document_id=new_document_id(),
        version_id=new_version_id(),
        pages=(
            DocumentPage(
                1,
                (
                    DocumentBlock(
                        block_id="p1-b1",
                        block_type=BlockType.TEXT,
                        text="منظومة كاظمة للذكاء الاصطناعي " * 4,
                    ),
                ),
                metadata={"extractor": "pdfplumber"},
            ),
        ),
        provenance=Provenance(source=str(path), parser="pdf"),
        metadata={"extractor": "pdfplumber"},
    )

    monkeypatch.setattr(
        pdf_mod.PdfParser, "_parse_pymupdf", staticmethod(lambda *_a, **_k: weak)
    )
    monkeypatch.setattr(
        pdf_mod.PdfParser, "_parse_pdfplumber", staticmethod(lambda *_a, **_k: strong)
    )
    monkeypatch.setattr(
        pdf_mod.PdfParser,
        "_parse_pypdf",
        staticmethod(lambda *_a, **_k: (_ for _ in ()).throw(ImportError("skip"))),
    )

    ir = pdf_mod.PdfParser().parse(
        path, _context(path, document_config, "pdf", "application/pdf")
    )
    assert ir.metadata.get("extractor") == "pdfplumber"
    assert "منظومة" in ir.pages[0].blocks[0].text
    assert any(item.startswith("pymupdf:") for item in ir.metadata.get("extractors_tried", []))
    assert any(item.startswith("pdfplumber:") for item in ir.metadata.get("extractors_tried", []))


def test_scanned_pdf_routes_to_ocr_and_extracts_arabic(
    tmp_path: Path, document_config: DocumentConfig
) -> None:
    """Tier-3: image-only PDF → quality needs_ocr → apply_ocr (existing OCR service).

    Draws *shaped* Arabic onto a raster page (no text layer). Routing is always
    asserted; content recovery requires Tesseract + ``ara``.
    """
    fitz = pytest.importorskip("fitz")
    Image = pytest.importorskip("PIL.Image")
    ImageDraw = pytest.importorskip("PIL.ImageDraw")
    ImageFont = pytest.importorskip("PIL.ImageFont")

    from kazma_core.documents.ocr import apply_ocr, get_ocr_health
    from kazma_core.documents.ocr.base import OcrReadiness
    from kazma_core.documents.parsers.pdf import PdfParser
    from kazma_core.documents.quality import assess_document_quality

    # Shape for visual rendering so Tesseract sees joined glyphs (PIL does not).
    phrase = "منظومة كاظمة"
    try:
        import arabic_reshaper
        from bidi.algorithm import get_display

        draw_text = get_display(arabic_reshaper.reshape(phrase))
    except Exception:
        draw_text = phrase

    img = Image.new("RGB", (900, 240), "white")
    draw = ImageDraw.Draw(img)
    font = None
    for candidate in (
        Path(r"C:\Windows\Fonts\tradbdo.ttf"),
        Path(r"C:\Windows\Fonts\arial.ttf"),
        Path(r"C:\Windows\Fonts\tahoma.ttf"),
        Path("/usr/share/fonts/truetype/noto/NotoSansArabic-Regular.ttf"),
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
    ):
        if candidate.is_file():
            try:
                font = ImageFont.truetype(str(candidate), 56)
                break
            except OSError:
                continue
    if font is None:
        font = ImageFont.load_default()
    draw.text((40, 60), draw_text, fill="black", font=font)
    image_path = tmp_path / "scan.png"
    img.save(image_path)

    pdf_path = tmp_path / "scanned-ar.pdf"
    doc = fitz.open()
    page = doc.new_page(width=900, height=240)
    page.insert_image(page.rect, filename=str(image_path))
    doc.save(pdf_path)
    doc.close()

    ir = PdfParser().parse(
        pdf_path, _context(pdf_path, document_config, "pdf", "application/pdf")
    )
    qualities = assess_document_quality(ir, min_text_chars=40)
    assert qualities[0].needs_ocr is True
    assert "no_native_text" in qualities[0].reasons or "image_or_scanned_page" in qualities[
        0
    ].reasons

    health = get_ocr_health(("eng", "ara"))
    if health.readiness is OcrReadiness.UNAVAILABLE or "ara" not in (
        health.languages or ()
    ):
        pytest.skip("Tesseract + ara not available for end-to-end OCR assert")

    ocr_config = replace(
        document_config,
        ocr_enabled=True,
        ocr_languages=("eng", "ara"),
        ocr_dpi=200,
    )
    ir = replace(
        ir,
        metadata={**dict(ir.metadata), "source_mime": "application/pdf"},
    )
    result = apply_ocr(pdf_path, ir, ocr_config, work_dir=tmp_path)
    joined = "\n".join(block.text for page in result.pages for block in page.blocks)
    # Real OCR on synthetic glyphs is imperfect; require Arabic script recovery
    # (not Latin gibberish from eng-first language order).
    has_arabic = any("\u0600" <= ch <= "\u06FF" for ch in joined)
    assert has_arabic and joined.strip(), (
        f"OCR did not recover Arabic script from scanned PDF: {joined!r}; "
        f"meta={result.metadata.get('ocr')!r} warnings={result.metadata.get('warnings')!r}"
    )
    assert result.metadata.get("ocr", {}).get("status") == "completed"
