"""Document Processor Native Skill — read, merge, split, inspect, convert, OCR, and manipulate documents.

Each tool lazily imports its heavy dependency and returns a friendly
install-hint string when the library is missing, so the skill always loads.
Output is written to ``kazma-data/documents/`` for generated/manipulated files.

Optional extras:
- ``pip install 'kazma[ocr]'`` for OCR (scanned PDFs, images).
- ``pip install 'kazma[convert]'`` for format conversion (HTML→PDF, MD→PDF).
"""

from __future__ import annotations

import csv
import json
import logging
import re
import time
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

DOC_DIR = Path("kazma-data/documents")

# Supported document suffixes (used by file_read delegation + attachment parsing)
DOCUMENT_SUFFIXES = frozenset({
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".csv",
    ".pptx", ".ppt", ".json", ".txt", ".md", ".log", ".rtf",
})

_MAX_OUTPUT_CHARS = 20_000  # generous cap (vs the old 8_000)


# ── Helpers ─────────────────────────────────────────────────────────────


def _slugify(text: str, max_len: int = 50) -> str:
    slug = re.sub(r"[^a-z0-9\s-]", "", text.lower().strip())
    slug = re.sub(r"[\s-]+", "-", slug).strip("-")
    return (slug or "document")[:max_len]


def _filename(name: str, ext: str) -> Path:
    DOC_DIR.mkdir(parents=True, exist_ok=True)
    ts = int(time.time())
    return DOC_DIR / f"{ts}_{_slugify(name)}.{ext}"


def _truncate(text: str, limit: int = _MAX_OUTPUT_CHARS) -> str:
    if len(text) > limit:
        return text[:limit] + f"\n[... truncated at {limit} chars ...]"
    return text


# ── read_document ───────────────────────────────────────────────────────


async def read_document(path: str) -> str:
    """Read and extract text from a document file.

    Auto-detects format by suffix and extracts text + tables + metadata.
    Supported: PDF, DOCX, XLSX, PPTX, CSV, JSON, TXT, MD.

    Args:
        path: Path to the document file.

    Returns:
        Extracted text content with page/sheet markers, or an error message.
    """
    p = Path(path).expanduser().resolve()
    if not p.exists():
        return f"Error: Document not found: {path}"
    if not p.is_file():
        return f"Error: Path is not a file: {path}"

    suffix = p.suffix.lower()

    try:
        if suffix == ".pdf":
            return _read_pdf(p)
        elif suffix == ".docx":
            return _read_docx(p)
        elif suffix in (".xlsx", ".xls"):
            return _read_excel(p)
        elif suffix == ".pptx":
            return _read_pptx(p)
        elif suffix == ".csv":
            return _read_csv(p)
        elif suffix == ".json":
            return _read_json(p)
        elif suffix in (".txt", ".md", ".log"):
            return _read_text(p)
        elif suffix == ".rtf":
            return _read_rtf(p)
        else:
            # Last resort: try as text
            return _read_text(p)
    except Exception as exc:
        logger.error("[document_processor] Error reading %s: %s", path, exc)
        return f"Error reading document: {exc}"


def _read_pdf(p: Path) -> str:
    """Extract text + tables from a PDF using pdfplumber (primary) or pypdf (fallback)."""
    # Primary: pdfplumber (better quality + table extraction)
    try:
        import pdfplumber

        lines: list[str] = []
        with pdfplumber.open(p) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                if text.strip():
                    lines.append(f"--- Page {i + 1} ---")
                    lines.append(text.strip())
                    lines.append("")

                # Extract tables if present
                try:
                    tables = page.extract_tables()
                    for j, table in enumerate(tables):
                        if table and len(table) > 1:
                            lines.append(f"  [Table {j + 1} on page {i + 1}]")
                            for row in table[:20]:  # cap at 20 rows per table
                                cells = [str(c) if c else "" for c in row]
                                lines.append("  | " + " | ".join(cells) + " |")
                            if len(table) > 20:
                                lines.append(f"  [... {len(table) - 20} more rows ...]")
                            lines.append("")
                except Exception:
                    pass  # table extraction is best-effort

        result = "\n".join(lines)
        if result.strip():
            return _truncate(result)

        # No text extracted — might be a scanned PDF
        return (
            "No text layer found in this PDF. It may be a scanned/image-only document. "
            "Install OCR support: pip install 'kazma[ocr]'"
        )
    except ImportError:
        pass
    except Exception as exc:
        logger.debug("[document_processor] pdfplumber failed: %s — trying pypdf", exc)

    # Fallback: pypdf
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(p))
        lines = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                lines.append(f"--- Page {i + 1} ---")
                lines.append(text.strip())
                lines.append("")
        result = "\n".join(lines)
        if result.strip():
            return _truncate(result)
        return (
            "No text layer found in this PDF. It may be a scanned/image-only document. "
            "Install OCR support: pip install 'kazma[ocr]'"
        )
    except ImportError:
        return "Error: No PDF library installed. Run: pip install pdfplumber pypdf"
    except Exception as exc:
        return f"Error parsing PDF: {exc}"


def _read_docx(p: Path) -> str:
    """Extract text from a DOCX file using python-docx."""
    try:
        from docx import Document
    except ImportError:
        return "Error: python-docx not installed. Run: pip install python-docx"

    doc = Document(str(p))
    lines: list[str] = []

    # Paragraphs (including headings)
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # Detect heading style for structure markers
            style_name = (para.style.name or "").lower()
            if "heading" in style_name:
                level = re.search(r"\d+", style_name)
                prefix = "#" * int(level.group()) if level else "#"
                lines.append(f"{prefix} {text}")
            else:
                lines.append(text)

    # Tables
    for i, table in enumerate(doc.tables):
        lines.append(f"\n[Table {i + 1}]")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")

    result = "\n".join(lines)
    return _truncate(result) if result else "Document appears to be empty."


def _read_excel(p: Path) -> str:
    """Extract data from XLSX/XLS using openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "Error: openpyxl not installed. Run: pip install openpyxl"

    wb = load_workbook(p, read_only=True, data_only=True)
    lines: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"--- Sheet: {sheet_name} ---")
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= 200:
                lines.append("[... truncated after 200 rows ...]")
                break
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                lines.append(" | ".join(cells))
        lines.append("")
    wb.close()
    return _truncate("\n".join(lines))


def _read_pptx(p: Path) -> str:
    """Extract text from PPTX slides using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        return "Error: python-pptx not installed. Run: pip install python-pptx"

    prs = Presentation(str(p))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides):
        lines.append(f"--- Slide {i + 1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                table = shape.table
                lines.append(f"  [Table on slide {i + 1}]")
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append("  | " + " | ".join(cells) + " |")
        lines.append("")
    return _truncate("\n".join(lines))


def _read_csv(p: Path) -> str:
    """Extract data from a CSV file."""
    lines: list[str] = []
    with open(p, "r", encoding="utf-8", errors="replace", newline="") as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            if i >= 200:
                lines.append("[... truncated after 200 rows ...]")
                break
            lines.append(", ".join(row))
    return _truncate("\n".join(lines))


def _read_json(p: Path) -> str:
    """Pretty-print a JSON file."""
    with open(p, "r", encoding="utf-8", errors="replace") as f:
        data = json.load(f)
    return _truncate(json.dumps(data, indent=2, ensure_ascii=False))


def _read_text(p: Path) -> str:
    """Read a plain-text file."""
    content = p.read_text(encoding="utf-8", errors="replace")
    return _truncate(content)


def _read_rtf(p: Path) -> str:
    """Strip RTF formatting to get plain text (best-effort)."""
    import re as _re

    raw = p.read_text(encoding="utf-8", errors="replace")
    # Remove RTF control words and groups
    text = _re.sub(r"\\[a-z]+-?\d+ ?|[@{}]|\\'[0-9a-f]{2}", "", raw, flags=_re.IGNORECASE)
    text = _re.sub(r"[{}]", "", text)
    text = _re.sub(r"\s+", " ", text).strip()
    return _truncate(text) if text else "RTF file appears to be empty."


# ── pdf_merge ───────────────────────────────────────────────────────────


async def pdf_merge(file_paths: list[str], output_name: str = "merged") -> str:
    """Merge multiple PDF files into a single PDF.

    Args:
        file_paths: List of PDF file paths to merge (in order).
        output_name: Name for the output file (without extension).

    Returns:
        Success message with the output path, or error.
    """
    try:
        from pypdf import PdfWriter
    except ImportError:
        return "Error: pypdf not installed. Run: pip install pypdf"

    if not file_paths:
        return "Error: No file paths provided."

    writer = PdfWriter()
    for fp in file_paths:
        p = Path(fp).expanduser().resolve()
        if not p.exists():
            return f"Error: File not found: {fp}"
        if p.suffix.lower() != ".pdf":
            return f"Error: Not a PDF file: {fp}"
        try:
            writer.append(str(p))
        except Exception as exc:
            return f"Error: Failed to append {fp}: {exc}"

    dest = _filename(output_name, "pdf")
    with open(dest, "wb") as f:
        writer.write(f)
    writer.close()

    return f"PDF merged successfully.\n  Files: {len(file_paths)}\n  Saved to: {dest}"


# ── pdf_split ───────────────────────────────────────────────────────────


async def pdf_split(
    file_path: str,
    start_page: int = 1,
    end_page: int = 0,
    output_name: str = "",
) -> str:
    """Extract specific pages from a PDF into a new file.

    Args:
        file_path: Path to the source PDF.
        start_page: First page to extract (1-indexed).
        end_page: Last page to extract (1-indexed, 0 = to end).
        output_name: Name for the output file (without extension).

    Returns:
        Success message with the output path, or error.
    """
    try:
        from pypdf import PdfReader, PdfWriter
    except ImportError:
        return "Error: pypdf not installed. Run: pip install pypdf"

    p = Path(file_path).expanduser().resolve()
    if not p.exists():
        return f"Error: File not found: {file_path}"

    reader = PdfReader(str(p))
    total = len(reader.pages)

    if start_page < 1:
        start_page = 1
    if end_page <= 0 or end_page > total:
        end_page = total
    if start_page > total:
        return f"Error: start_page ({start_page}) exceeds total pages ({total})."

    writer = PdfWriter()
    for i in range(start_page - 1, min(end_page, total)):
        writer.add_page(reader.pages[i])

    name = output_name or f"{p.stem}_pages_{start_page}-{end_page}"
    dest = _filename(name, "pdf")
    with open(dest, "wb") as f:
        writer.write(f)
    writer.close()

    return (
        f"PDF split successfully.\n"
        f"  Source: {p.name} ({total} pages)\n"
        f"  Extracted: pages {start_page}-{end_page}\n"
        f"  Saved to: {dest}"
    )


# ── pdf_info ────────────────────────────────────────────────────────────


async def pdf_info(file_path: str) -> str:
    """Extract metadata from a PDF file.

    Args:
        file_path: Path to the PDF file.

    Returns:
        Formatted metadata string (page count, title, author, dates, etc.).
    """
    try:
        from pypdf import PdfReader
    except ImportError:
        return "Error: pypdf not installed. Run: pip install pypdf"

    p = Path(file_path).expanduser().resolve()
    if not p.exists():
        return f"Error: File not found: {file_path}"

    reader = PdfReader(str(p))
    total = len(reader.pages)
    meta = reader.metadata or {}

    # Check for text layer
    has_text = False
    for page in reader.pages[:5]:  # sample first 5 pages
        if (page.extract_text() or "").strip():
            has_text = True
            break

    lines = [
        f"PDF Metadata: {p.name}",
        f"  Pages: {total}",
        f"  Title: {meta.get('/Title', '(none)')}",
        f"  Author: {meta.get('/Author', '(none)')}",
        f"  Subject: {meta.get('/Subject', '(none)')}",
        f"  Creator: {meta.get('/Creator', '(none)')}",
        f"  Created: {meta.get('/CreationDate', '(unknown)')}",
        f"  Modified: {meta.get('/ModDate', '(unknown)')}",
        f"  Text layer: {'Yes' if has_text else 'No (likely scanned/image-only)'}",
    ]
    if not has_text:
        lines.append("  ⚠ No extractable text. Install OCR: pip install 'kazma[ocr]'")

    return "\n".join(lines)


# ── ocr_document (optional: pip install 'kazma[ocr]') ───────────────────


async def ocr_document(path: str, lang: str = "eng") -> str:
    """OCR a scanned PDF or image to extract text.

    Requires the ``[ocr]`` extra: ``pip install 'kazma[ocr]'`` plus a
    system install of Tesseract OCR (``apt install tesseract-ocr`` on
    Linux, ``brew install tesseract`` on macOS, or the Windows installer
    from https://github.com/UB-Mannheim/tesseract/wiki).

    Args:
        path: Path to a PDF or image file (PNG/JPEG/TIFF).
        lang: Tesseract language code (e.g. ``"eng"``, ``"ara"``, ``"eng+ara"``).

    Returns:
        Extracted text, or an install-hint error if OCR deps are missing.
    """
    p = Path(path).expanduser().resolve()
    if not p.exists():
        return f"Error: File not found: {path}"

    suffix = p.suffix.lower()
    try:
        import pytesseract
    except ImportError:
        return (
            "Error: OCR libraries not installed. Run: pip install 'kazma[ocr]' "
            "and install Tesseract OCR (apt install tesseract-ocr / brew install tesseract "
            "/ Windows installer from UB-Mannheim/tesseract wiki)."
        )

    try:
        if suffix == ".pdf":
            # Convert PDF pages to images, then OCR each
            try:
                from pdf2image import convert_from_path
            except ImportError:
                return "Error: pdf2image not installed. Run: pip install 'kazma[ocr]'"

            try:
                from pypdf import PdfReader

                total_pages = len(PdfReader(str(p)).pages)
            except Exception:
                total_pages = 10  # fallback

            images = convert_from_path(str(p), dpi=200)
            lines: list[str] = []
            for i, img in enumerate(images):
                text = pytesseract.image_to_string(img, lang=lang)
                if text.strip():
                    lines.append(f"--- Page {i + 1} ---")
                    lines.append(text.strip())
                    lines.append("")
            return _truncate("\n".join(lines)) if lines else "No text detected via OCR."

        else:
            # Image file (PNG/JPEG/TIFF)
            from PIL import Image

            img = Image.open(str(p))
            text = pytesseract.image_to_string(img, lang=lang)
            return _truncate(text) if text.strip() else "No text detected in image."

    except Exception as exc:
        return f"Error during OCR: {exc}"


# ── convert_document (optional: pip install 'kazma[convert]') ───────────


async def convert_document(
    file_path: str,
    target_format: str,
    output_name: str = "",
) -> str:
    """Convert a document between formats.

    Supported conversions (requires optional libs):
    - HTML → PDF (``weasyprint``)
    - Markdown → PDF (``markdown`` + ``weasyprint``)
    - Markdown → HTML (``markdown`` — built-in)
    - Markdown → DOCX (``python-docx`` + ``markdown``)

    Args:
        file_path: Source file path.
        target_format: Target format (``pdf``, ``html``, ``docx``).
        output_name: Output filename (without extension).

    Returns:
        Success message with output path, or error.
    """
    p = Path(file_path).expanduser().resolve()
    if not p.exists():
        return f"Error: File not found: {file_path}"

    src_suffix = p.suffix.lower()
    target = target_format.lower().lstrip(".")
    name = output_name or p.stem
    src_text = p.read_text(encoding="utf-8", errors="replace")

    try:
        # ── Markdown → HTML ──
        if src_suffix in (".md", ".markdown") and target == "html":
            import markdown as md

            html = md.markdown(src_text, extensions=["tables", "fenced_code"])
            dest = _filename(name, "html")
            dest.write_text(f"<html><body>{html}</body></html>", encoding="utf-8")
            return f"Converted successfully.\n  {src_suffix} → .html\n  Saved to: {dest}"

        # ── Markdown/HTML → PDF ──
        if src_suffix in (".md", ".markdown", ".html", ".htm") and target == "pdf":
            try:
                from weasyprint import HTML
            except ImportError:
                return "Error: weasyprint not installed. Run: pip install 'kazma[convert]'"

            if src_suffix in (".md", ".markdown"):
                import markdown as md

                html_content = md.markdown(src_text, extensions=["tables", "fenced_code"])
                full_html = f"<html><body style='font-family: sans-serif; font-size: 12pt;'>{html_content}</body></html>"
            else:
                full_html = src_text

            dest = _filename(name, "pdf")
            HTML(string=full_html).write_pdf(str(dest))
            return f"Converted successfully.\n  {src_suffix} → .pdf\n  Saved to: {dest}"

        # ── Markdown → DOCX ──
        if src_suffix in (".md", ".markdown") and target == "docx":
            from docx import Document as DocxDoc

            doc = DocxDoc()
            for line in src_text.splitlines():
                stripped = line.strip()
                if stripped.startswith("#"):
                    level = min(stripped.count("#"), 3)
                    text = stripped.lstrip("#").strip()
                    doc.add_heading(text, level=level)
                elif stripped:
                    doc.add_paragraph(stripped)
            dest = _filename(name, "docx")
            doc.save(str(dest))
            return f"Converted successfully.\n  .md → .docx\n  Saved to: {dest}"

        return f"Error: Unsupported conversion {src_suffix} → .{target}"

    except ImportError as exc:
        return f"Error: Missing library for this conversion: {exc}. Run: pip install 'kazma[convert]'"
    except Exception as exc:
        return f"Error during conversion: {exc}"


# ── pdf_fill_form ───────────────────────────────────────────────────────


async def pdf_fill_form(
    file_path: str,
    fields: dict[str, str],
    output_name: str = "",
) -> str:
    """Fill AcroForm fields in a PDF.

    Args:
        file_path: Path to the source PDF with form fields.
        fields: Dict mapping field names to values.
        output_name: Output filename (without extension).

    Returns:
        Success message with output path, or error.
    """
    try:
        from pypdf import PdfReader, PdfWriter
    except ImportError:
        return "Error: pypdf not installed. Run: pip install pypdf"

    p = Path(file_path).expanduser().resolve()
    if not p.exists():
        return f"Error: File not found: {file_path}"

    reader = PdfReader(str(p))
    writer = PdfWriter()
    writer.append(reader)

    filled = 0
    for page in writer.pages:
        writer.update_page_form_field_values(page, fields)
        filled += 1

    # Check if form fields exist
    field_count = 0
    for page in writer.pages:
        if "/Annots" in page:
            for annot_ref in page["/Annots"]:
                annot = annot_ref.get_object()
                if annot.get("/Subtype") == "/Widget":
                    field_count += 1

    if field_count == 0:
        return f"Warning: No AcroForm fields detected in this PDF. It may use XFA forms (not supported) or have no form fields."

    name = output_name or f"{p.stem}_filled"
    dest = _filename(name, "pdf")
    with open(dest, "wb") as f:
        writer.write(f)
    writer.close()

    return f"PDF form filled.\n  Fields provided: {len(fields)}\n  Pages: {filled}\n  Saved to: {dest}"


# ── pdf_redact ──────────────────────────────────────────────────────────


async def pdf_redact(
    file_path: str,
    terms: list[str],
    output_name: str = "",
) -> str:
    """Redact text from a PDF by covering matched terms with black rectangles.

    Args:
        file_path: Path to the source PDF.
        terms: List of text strings to redact (case-insensitive).
        output_name: Output filename (without extension).

    Returns:
        Success message with output path + redaction count.
    """
    try:
        import pdfplumber
        from pypdf import PdfReader, PdfWriter
        from reportlab.pdfgen import canvas
        from reportlab.lib.colors import black
    except ImportError as exc:
        return f"Error: Missing library: {exc}. Run: pip install pdfplumber pypdf reportlab"

    p = Path(file_path).expanduser().resolve()
    if not p.exists():
        return f"Error: File not found: {file_path}"

    terms_lower = [t.lower() for t in terms]
    total_redacted = 0

    reader = PdfReader(str(p))
    writer = PdfWriter()

    for page_idx, page in enumerate(reader.pages):
        # Find text positions to redact using pdfplumber
        import io

        from reportlab.lib.pagesizes import letter as _letter

        page_w = float(page.mediabox.width)
        page_h = float(page.mediabox.height)

        # Create overlay with black rectangles
        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=(page_w, page_h))

        try:
            with pdfplumber.open(p) as pdf:
                pg = pdf.pages[page_idx]
                for word in pg.extract_words():
                    word_text = word["text"].lower()
                    for term in terms_lower:
                        if term in word_text:
                            # Draw black rectangle over the word
                            c.setFillColor(black)
                            c.rect(
                                word["x0"],
                                page_h - word["bottom"],
                                word["x1"] - word["x0"],
                                word["bottom"] - word["top"],
                                fill=1,
                                stroke=0,
                            )
                            total_redacted += 1
                            break
        except Exception:
            pass  # best-effort — if extraction fails, skip redaction on this page

        c.save()
        buf.seek(0)

        # Merge overlay onto page
        from pypdf import PdfReader as _PR

        overlay_reader = _PR(buf)
        page.merge_page(overlay_reader.pages[0])
        writer.add_page(page)

    name = output_name or f"{p.stem}_redacted"
    dest = _filename(name, "pdf")
    with open(dest, "wb") as f:
        writer.write(f)
    writer.close()

    return (
        f"PDF redacted successfully.\n"
        f"  Terms: {', '.join(terms)}\n"
        f"  Instances redacted: {total_redacted}\n"
        f"  Saved to: {dest}"
    )


# ── generate_pptx ───────────────────────────────────────────────────────


async def generate_pptx(
    title: str,
    slides: list[dict[str, Any]],
) -> str:
    """Generate a PowerPoint presentation from a title and slide data.

    Requires ``python-pptx`` (included in the ``[document]`` extra).

    Args:
        title: Presentation title (used for the title slide).
        slides: List of slide dicts. Each slide supports:
            - ``heading`` (str): Slide title.
            - ``body`` (str): Bullet text (newline-separated for multiple bullets).
            - ``bullets`` (list[str]): Explicit bullet list (alternative to body).
            - ``layout`` (str): "title" | "content" (default "content").

    Returns:
        Success message with the saved file path.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return "Error: python-pptx not installed. Run: pip install python-pptx"

    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = ""

    # Content slides
    content_layout = prs.slide_layouts[1]  # Title + Content
    for s in slides:
        heading = s.get("heading", "")
        body = s.get("body", "")
        bullets = s.get("bullets")
        layout_name = s.get("layout", "content")

        if layout_name == "title":
            slide_layout = prs.slide_layouts[0]
        else:
            slide_layout = content_layout

        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = heading

        # Add body text as bullets
        if bullets:
            text_frame = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
            if text_frame:
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        text_frame.text = str(bullet)
                    else:
                        p = text_frame.add_paragraph()
                        p.text = str(bullet)
                        p.level = 0
        elif body:
            text_frame = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
            if text_frame:
                lines = body.split("\n")
                for i, line in enumerate(lines):
                    line = line.strip()
                    if not line:
                        continue
                    if i == 0 and not text_frame.text:
                        text_frame.text = line
                    else:
                        p = text_frame.add_paragraph()
                        p.text = line
                        # Support indentation via leading spaces
                        p.level = 1 if line.startswith("  ") else 0

    dest = _filename(title, "pptx")
    prs.save(str(dest))

    return f"PowerPoint generated successfully.\n  Title: {title}\n  Slides: {len(slides) + 1}\n  Saved to: {dest}"
