"""JSON-file protocol worker for document generation and conversion."""

from __future__ import annotations

import hashlib
import html
import json
import os
import re
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Any

from .resources import validate_restricted_render_resources

_PROTOCOL_VERSION = 1
_MIME = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "html": "text/html",
    "markdown": "text/markdown",
    "md": "text/markdown",
}


class WorkerError(RuntimeError):
    def __init__(self, code: str, message: str) -> None:
        super().__init__(message)
        self.code = code
        self.safe_message = message


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        while chunk := handle.read(1024 * 1024):
            digest.update(chunk)
    return digest.hexdigest()


def _write(path: Path, value: dict[str, Any]) -> None:
    payload = json.dumps(
        value, ensure_ascii=False, allow_nan=False, sort_keys=True, separators=(",", ":")
    ).encode("utf-8")
    temporary = path.with_suffix(".writing")
    temporary.write_bytes(payload)
    os.replace(temporary, path)


def _request(path: Path) -> dict[str, Any]:
    value = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(value, dict) or value.get("protocol_version") != _PROTOCOL_VERSION:
        raise WorkerError("invalid_request", "Invalid renderer worker request")
    return value


def _verify_assets(request: dict[str, Any], work_dir: Path) -> None:
    records = request.get("approved_assets", [])
    if not isinstance(records, list):
        raise WorkerError("invalid_request", "Invalid approved render assets")
    expected: set[str] = set()
    assets = work_dir / "assets"
    for record in records:
        if not isinstance(record, dict) or set(record) != {"name", "sha256"}:
            raise WorkerError("invalid_request", "Invalid approved render asset record")
        name = str(record["name"])
        if Path(name).name != name or name in expected:
            raise WorkerError("invalid_request", "Invalid approved render asset name")
        path = assets / name
        if not path.is_file() or _sha256(path) != record["sha256"]:
            raise WorkerError("document_changed", "An approved render asset changed")
        expected.add(name)
    actual = {path.name for path in assets.iterdir() if path.is_file()} if assets.is_dir() else set()
    if actual != expected:
        raise WorkerError("invalid_request", "Unexpected files in render asset directory")


def _validate_payload_limits(request: dict[str, Any]) -> None:
    payload = request.get("payload")
    limits = request.get("limits")
    if not isinstance(payload, dict) or not isinstance(limits, dict):
        raise WorkerError("invalid_request", "Invalid renderer payload limits")
    sheets = payload.get("sheets", [])
    if isinstance(sheets, list):
        if len(sheets) > int(limits["max_sheets"]):
            raise WorkerError("document_limit_exceeded", "Workbook exceeds sheet limit")
        cells = 0
        for sheet in sheets:
            rows = sheet.get("rows", []) if isinstance(sheet, dict) else []
            if isinstance(rows, list) and len(rows) > int(limits["max_rows_per_sheet"]):
                raise WorkerError("document_limit_exceeded", "Workbook exceeds row limit")
            if isinstance(rows, list):
                cells += sum(len(row) for row in rows if isinstance(row, list))
        if cells > int(limits["max_cells"]):
            raise WorkerError("document_limit_exceeded", "Workbook exceeds cell limit")
    slides = payload.get("slides", [])
    if isinstance(slides, list) and len(slides) + 1 > int(limits["max_slides"]):
        raise WorkerError("document_limit_exceeded", "Presentation exceeds slide limit")
    images = payload.get("images", [])
    if isinstance(images, list) and len(images) > int(limits["max_images"]):
        raise WorkerError("document_limit_exceeded", "Document exceeds image limit")


def _sections(value: object) -> list[dict[str, str]]:
    if not isinstance(value, list):
        return []
    result: list[dict[str, str]] = []
    for item in value:
        if isinstance(item, dict):
            result.append(
                {"heading": str(item.get("heading", "")), "body": str(item.get("body", ""))}
            )
    return result


def _markdown(payload: dict[str, Any]) -> str:
    lines = [f"# {payload.get('title', 'Document')}", ""]
    if payload.get("toc"):
        lines.extend(("## Contents", ""))
        for section in _sections(payload.get("sections")):
            heading = section["heading"].lstrip("#").strip()
            if heading:
                anchor = re.sub(r"[^\w\u0600-\u06ff -]", "", heading.lower()).replace(" ", "-")
                lines.append(f"- [{heading}](#{anchor})")
        lines.append("")
    for section in _sections(payload.get("sections")):
        if section["heading"]:
            lines.extend((f"## {section['heading'].lstrip('#').strip()}", ""))
        if section["body"]:
            lines.extend((section["body"], ""))
    citations = payload.get("citations")
    if isinstance(citations, list) and citations:
        lines.extend(("## References", ""))
        lines.extend(f"{index}. {item}" for index, item in enumerate(citations, 1))
    return "\n".join(lines).rstrip() + "\n"


def _markdown_html(text: str) -> str:
    try:
        import markdown

        body = markdown.markdown(text, extensions=["tables", "fenced_code", "toc"])
    except ImportError:
        body = "<pre>" + html.escape(text) + "</pre>"
    return (
        "<!doctype html><html><head><meta charset=\"utf-8\">"
        "<style>body{font-family:sans-serif;unicode-bidi:plaintext}"
        "[dir=rtl]{direction:rtl} table{border-collapse:collapse}</style></head>"
        f"<body>{body}</body></html>"
    )


def _safe_html(text: str) -> None:
    try:
        validate_restricted_render_resources(text)
    except Exception as exc:
        raise WorkerError(
            "external_resource_denied",
            "External or unapproved local resources are forbidden during rendering",
        ) from exc


def _font_paths() -> tuple[Path | None, Path | None]:
    """Prefer fonts with solid Arabic coverage (glyphs + metrics)."""
    candidates = (
        # Windows — Calibri first for exact user preference
        (Path("C:/Windows/Fonts/calibri.ttf"), Path("C:/Windows/Fonts/calibrib.ttf")),
        (Path("C:/Windows/Fonts/arial.ttf"), Path("C:/Windows/Fonts/arialbd.ttf")),
        (Path("C:/Windows/Fonts/tahoma.ttf"), Path("C:/Windows/Fonts/tahomabd.ttf")),
        (Path("C:/Windows/Fonts/trado.ttf"), Path("C:/Windows/Fonts/trado.ttf")),
        (
            Path("C:/Windows/Fonts/NotoSansArabic-Regular.ttf"),
            Path("C:/Windows/Fonts/NotoSansArabic-Bold.ttf"),
        ),
        (
            Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
            Path("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
        ),
        (
            Path("/usr/share/fonts/truetype/noto/NotoSansArabic-Regular.ttf"),
            Path("/usr/share/fonts/truetype/noto/NotoSansArabic-Bold.ttf"),
        ),
        (
            Path("/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf"),
            Path("/usr/share/fonts/truetype/noto/NotoSans-Bold.ttf"),
        ),
        (
            Path("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"),
            Path("/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"),
        ),
    )
    return next(((regular, bold) for regular, bold in candidates if regular.is_file()), (None, None))


def _generate_pdf(output: Path, payload: dict[str, Any], warnings: list[str]) -> None:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT, TA_RIGHT
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import (
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    from kazma_core.documents.rich_render import (
        inline_markdown_to_reportlab,
        is_arabic_dominant,
        pdf_flowables_from_body,
        shape_for_pdf,
    )

    regular, bold = _font_paths()
    font = "Helvetica"
    bold_font = "Helvetica-Bold"
    if regular:
        pdfmetrics.registerFont(TTFont("KazmaUnicode", str(regular)))
        font = "KazmaUnicode"
        if bold and bold.is_file():
            pdfmetrics.registerFont(TTFont("KazmaUnicodeBold", str(bold)))
            bold_font = "KazmaUnicodeBold"
        else:
            bold_font = font
    else:
        warnings.append("Unicode font unavailable; PDF uses a limited deterministic fallback")
    style = payload.get("style") if isinstance(payload.get("style"), dict) else {}

    def size(name: str, default: float, low: float, high: float) -> float:
        try:
            return min(high, max(low, float(style.get(name, default))))
        except (TypeError, ValueError):
            warnings.append(f"Invalid {name} style token; deterministic default applied")
            return default

    # Detect document language from title + section text
    sections = _sections(payload.get("sections"))
    sample_parts = [str(payload.get("title", ""))]
    for item in sections:
        sample_parts.append(item.get("heading", ""))
        sample_parts.append((item.get("body") or "")[:1500])
    sample = "\n".join(sample_parts)
    rtl = is_arabic_dominant(sample)
    # Explicit payload override: lang=ar|en or rtl=true
    lang = str(payload.get("lang") or payload.get("language") or "").strip().lower()
    if lang in ("ar", "arabic", "rtl"):
        rtl = True
    elif lang in ("en", "english", "ltr"):
        rtl = False
    if payload.get("rtl") is True:
        rtl = True
    if payload.get("rtl") is False:
        rtl = False

    # ── Unified theme (EN == AR visual design; only dir/shape differ) ──
    from kazma_core.documents.style_theme import (
        THEME,
        localized_chrome,
        theme_colors_reportlab,
    )

    chrome = localized_chrome(rtl=rtl)
    th = theme_colors_reportlab()
    # Justified body in both languages. After shape_for_pdf() (reshape+get_display),
    # ReportLab must draw the *visual* string LTR — do NOT set wordWrap="RTL".
    shape_ar = rtl or is_arabic_dominant(sample)
    wrap = "LTR" if not rtl else "CJK"
    # Headings: same edge alignment for both — start-of-reading-direction
    align = TA_RIGHT if rtl else TA_LEFT
    body_align = TA_JUSTIFY

    title_size = size("title_font_size", float(THEME["title_size"]), 10, 36)
    heading_size = size("heading_font_size", float(THEME["h2_size"]), 8, 28)
    body_size = size("body_font_size", float(THEME["body_size"]), 6, 18)
    accent = th["accent"]
    if style.get("accent_color"):
        try:
            accent = colors.HexColor(str(style["accent_color"]))
        except Exception:
            pass
    heading_fill = th["heading_fill"]
    body_color = th["body"]

    # On-bar heading text is always white (same EN/AR)
    title_style = ParagraphStyle(
        "KazmaTitle",
        fontName=bold_font,
        fontSize=title_size,
        leading=title_size * 1.35,
        textColor=th["heading_text"],
        alignment=align,
        wordWrap=wrap,
        spaceBefore=0,
        spaceAfter=0,
    )
    h1_bar = ParagraphStyle(
        "KazmaH1Bar",
        fontName=bold_font,
        fontSize=float(THEME["h1_size"]),
        leading=float(THEME["h1_size"]) * 1.35,
        textColor=th["heading_text"],
        alignment=align,
        wordWrap=wrap,
        spaceBefore=0,
        spaceAfter=0,
    )
    h2_bar = ParagraphStyle(
        "KazmaH2Bar",
        fontName=bold_font,
        fontSize=heading_size,
        leading=heading_size * 1.35,
        textColor=th["heading_text"],
        alignment=align,
        wordWrap=wrap,
        spaceBefore=0,
        spaceAfter=0,
    )
    h3_style = ParagraphStyle(
        "KazmaH3",
        fontName=bold_font,
        fontSize=float(THEME["h3_size"]),
        leading=float(THEME["h3_size"]) * 1.4,
        spaceBefore=10,
        spaceAfter=4,
        textColor=th["heading"],
        alignment=align,
        wordWrap=wrap,
    )
    body_style = ParagraphStyle(
        "KazmaBody",
        fontName=font,
        fontSize=body_size,
        leading=body_size * float(THEME["line_height"]),
        textColor=body_color,
        alignment=body_align,
        wordWrap=wrap,
        spaceAfter=8,
        firstLineIndent=0,
    )
    # Body paragraphs: ragged-right RTL for Arabic (v4 layout), full-column
    # justified for English. body_style (TA_JUSTIFY) stays parent/fallback.
    body_para_style = ParagraphStyle(
        "KazmaBodyPara",
        parent=body_style,
        alignment=(TA_RIGHT if rtl else TA_JUSTIFY),
    )
    # Lists: same indent both sides so EN/AR feel symmetric
    bullet_style = ParagraphStyle(
        "KazmaBullet",
        parent=body_style,
        leftIndent=16,
        rightIndent=16,
        bulletIndent=0,
        alignment=align,
        spaceAfter=4,
    )
    number_style = ParagraphStyle("KazmaNumber", parent=bullet_style)
    quote_style = ParagraphStyle(
        "KazmaQuote",
        parent=body_style,
        textColor=th["quote"],
        leftIndent=12,
        rightIndent=12,
        backColor=th["bg_alt"],
        spaceBefore=4,
        spaceAfter=8,
    )
    code_style = ParagraphStyle(
        "KazmaCode",
        fontName=font,
        fontSize=8.5,
        leading=11,
        textColor=th["accent"],
        backColor=th["code_bg"],
        alignment=TA_LEFT,
        wordWrap="CJK",
        leftIndent=6,
        rightIndent=6,
        spaceBefore=4,
        spaceAfter=8,
    )
    rich_styles = {
        "body": body_style,
        "body_para": body_para_style,
        "h1": h1_bar,
        "h2": h2_bar,
        "h3": h3_style,
        "bullet": bullet_style,
        "number": number_style,
        "quote": quote_style,
        "code": code_style,
        "heading_fill": heading_fill,
    }

    def _bar(text_html: str, para_style: Any, *, fill: Any | None = None) -> Any:
        """Full-width filled heading bar — identical chrome for EN and AR."""
        para = Paragraph(text_html, para_style)
        fill_c = fill if fill is not None else heading_fill
        tbl = Table([[para]], colWidths=["*"])
        tbl.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, -1), fill_c),
                    ("LEFTPADDING", (0, 0), (-1, -1), 10),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 10),
                    ("TOPPADDING", (0, 0), (-1, -1), 8),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ]
            )
        )
        return tbl

    header = str(payload.get("header") or chrome["brand"])
    footer = str(payload.get("footer") or chrome["brand"])
    page_numbers = bool(payload.get("page_numbers", True))

    def decorate(canvas: Any, document: Any) -> None:
        canvas.saveState()
        canvas.setFont(font, 8)
        canvas.setFillColor(th["muted"])
        canvas.setStrokeColor(th["border"])
        canvas.setLineWidth(0.8)
        # Top brand rule
        canvas.line(
            document.leftMargin,
            A4[1] - 30,
            A4[0] - document.rightMargin,
            A4[1] - 30,
        )
        # Bottom rule
        canvas.line(
            document.leftMargin,
            36,
            A4[0] - document.rightMargin,
            36,
        )
        hdr = shape_for_pdf(header) if shape_ar else header
        ftr = shape_for_pdf(footer) if shape_ar else footer
        if rtl:
            canvas.drawRightString(A4[0] - document.rightMargin, A4[1] - 22, hdr)
            canvas.drawRightString(A4[0] - document.rightMargin, 22, ftr)
        else:
            canvas.drawString(document.leftMargin, A4[1] - 22, hdr)
            canvas.drawString(document.leftMargin, 22, ftr)
        if page_numbers:
            label = chrome["page_fmt"].format(n=document.page)
            if shape_ar:
                label = shape_for_pdf(label)
            canvas.drawCentredString(A4[0] / 2, 22, label)
        canvas.restoreState()

    title_raw = str(payload.get("title", "Document"))
    title_html = inline_markdown_to_reportlab(title_raw, shape_arabic=shape_ar)
    # Title bar uses accent (darker) — same card treatment EN/AR
    story: list[Any] = [
        _bar(title_html, title_style, fill=accent),
        Spacer(1, 14),
    ]

    if payload.get("toc"):
        story.append(
            _bar(
                inline_markdown_to_reportlab(chrome["toc"], shape_arabic=shape_ar),
                h2_bar,
            )
        )
        story.append(Spacer(1, 6))
        for index, item in enumerate(sections, 1):
            if item["heading"]:
                line = f"{index}. {item['heading']}"
                story.append(
                    Paragraph(
                        inline_markdown_to_reportlab(line, shape_arabic=shape_ar),
                        body_style,
                    )
                )
        story.append(PageBreak())

    for item in sections:
        if item["heading"]:
            story.append(
                _bar(
                    inline_markdown_to_reportlab(
                        item["heading"].lstrip("#").strip(),
                        shape_arabic=shape_ar,
                    ),
                    h1_bar,
                )
            )
            story.append(Spacer(1, 8))
        body = item.get("body") or ""
        if body.strip():
            col_width = float(A4[0]) - 2 * float(THEME.get("page_margin", 54))
            story.extend(
                pdf_flowables_from_body(
                    body,
                    styles=rich_styles,
                    shape_arabic=shape_ar,
                    Spacer=Spacer,
                    Paragraph=Paragraph,
                    colors=colors,
                    Table=Table,
                    TableStyle=TableStyle,
                    font_name=font,
                    bold_font_name=bold_font,
                    col_width=col_width,
                    font_size=body_size,
                )
            )

    tables = payload.get("tables")
    if isinstance(tables, list):
        for value in tables:
            if not isinstance(value, dict):
                continue
            heading = str(value.get("heading", ""))
            headers = value.get("headers")
            rows = value.get("rows")
            if heading:
                story.append(
                    _bar(
                        inline_markdown_to_reportlab(heading, shape_arabic=shape_ar),
                        h2_bar,
                    )
                )
                story.append(Spacer(1, 6))
            if isinstance(headers, list) and isinstance(rows, list) and headers:
                def _cell(val: object) -> str:
                    s = str(val)
                    return shape_for_pdf(s) if shape_ar else s

                data = [
                    [_cell(cell) for cell in headers],
                    *(
                        [_cell(cell) for cell in row]
                        for row in rows
                        if isinstance(row, list)
                    ),
                ]
                table = Table(data, repeatRows=1)
                table.setStyle(
                    TableStyle(
                        (
                            ("FONTNAME", (0, 0), (-1, -1), font),
                            ("FONTNAME", (0, 0), (-1, 0), bold_font),
                            ("BACKGROUND", (0, 0), (-1, 0), th["table_header_bg"]),
                            ("TEXTCOLOR", (0, 0), (-1, 0), th["table_header_fg"]),
                            ("BACKGROUND", (0, 1), (-1, -1), th["table_row_bg"]),
                            ("TEXTCOLOR", (0, 1), (-1, -1), body_color),
                            ("GRID", (0, 0), (-1, -1), 0.5, th["table_grid"]),
                            ("VALIGN", (0, 0), (-1, -1), "TOP"),
                            ("ALIGN", (0, 0), (-1, -1), "RIGHT" if rtl else "LEFT"),
                            ("TOPPADDING", (0, 0), (-1, -1), 6),
                            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                            ("LEFTPADDING", (0, 0), (-1, -1), 8),
                            ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                        )
                    )
                )
                story.extend((table, Spacer(1, 10)))
    citations = payload.get("citations")
    if isinstance(citations, list) and citations:
        story.append(
            _bar(
                inline_markdown_to_reportlab(
                    chrome["references"], shape_arabic=shape_ar
                ),
                h2_bar,
            )
        )
        story.append(Spacer(1, 6))
        for index, item in enumerate(citations, 1):
            line = f"{index}. {item}"
            story.append(
                Paragraph(
                    inline_markdown_to_reportlab(str(line), shape_arabic=shape_ar),
                    body_style,
                )
            )
    if payload.get("images"):
        warnings.append(
            "Images were omitted because generation accepts no unapproved filesystem resources"
        )
    if shape_ar:
        try:
            import arabic_reshaper  # noqa: F401
            from bidi.algorithm import get_display  # noqa: F401
        except ImportError:
            warnings.append(
                "arabic_reshaper/python-bidi not installed — Arabic letters may appear "
                "disconnected or reversed in PDF"
            )

    margin = float(THEME.get("page_margin", 54))
    SimpleDocTemplate(
        str(output),
        pagesize=A4,
        leftMargin=margin,
        rightMargin=margin,
        topMargin=margin + 4,
        bottomMargin=margin,
    ).build(story, onFirstPage=decorate, onLaterPages=decorate)


def _generate_docx(output: Path, payload: dict[str, Any]) -> None:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor

    from docx.shared import Cm

    from kazma_core.documents.rich_render import (
        docx_add_table,
        docx_apply_document_rtl,
        docx_force_justify,
        docx_heading_bar,
        docx_set_rtl_paragraph,
        docx_write_rich_body,
        is_arabic_dominant,
        try_parse_pipe_table_blob,
    )
    from kazma_core.documents.style_theme import THEME, localized_chrome

    document = Document()
    sections = _sections(payload.get("sections"))
    sample_parts = [str(payload.get("title", ""))]
    for item in sections:
        sample_parts.append(item.get("heading", ""))
        sample_parts.append((item.get("body") or "")[:1500])
    sample = "\n".join(sample_parts)
    rtl = is_arabic_dominant(sample)
    lang = str(payload.get("lang") or payload.get("language") or "").strip().lower()
    if lang in ("ar", "arabic", "rtl"):
        rtl = True
    elif lang in ("en", "english", "ltr"):
        rtl = False
    if payload.get("rtl") is True:
        rtl = True
    if payload.get("rtl") is False:
        rtl = False

    chrome = localized_chrome(rtl=rtl)
    fill_title = str(THEME["accent"]).lstrip("#")
    fill_h = str(THEME["heading_fill"]).lstrip("#")

    # Page setup — closer to PDF margins (not Word's 1.25" default)
    for section in document.sections:
        section.top_margin = Cm(1.8)
        section.bottom_margin = Cm(1.8)
        section.left_margin = Cm(1.8)
        section.right_margin = Cm(1.8)

    # Normal style: always justify (EN + AR); RTL via bidi on each para
    try:
        normal = document.styles["Normal"]
        normal.font.name = "Arial"
        normal.font.size = Pt(float(THEME["body_size"]))
        normal.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        normal.paragraph_format.space_after = Pt(8)
        normal.paragraph_format.line_spacing = float(THEME["line_height"])
    except Exception:
        pass

    # Title as full-width filled bar (PDF parity — works in Telegram/Word)
    docx_heading_bar(
        document,
        str(payload.get("title", "Document")),
        level=0,
        rtl=rtl,
        fill_hex=fill_title,
    )

    header = str(payload.get("header") or chrome["brand"])
    footer = str(payload.get("footer") or chrome["brand"])
    for section in document.sections:
        hp = section.header.paragraphs[0]
        hp.text = header
        if rtl:
            docx_set_rtl_paragraph(hp, justify=False)
        fp = section.footer.paragraphs[0]
        fp.text = footer
        if rtl:
            docx_set_rtl_paragraph(fp, justify=False)

    if payload.get("toc"):
        docx_heading_bar(document, chrome["toc"], level=1, rtl=rtl, fill_hex=fill_h)
        for index, item in enumerate(sections, 1):
            if item["heading"]:
                p = document.add_paragraph(f"{index}. {item['heading']}")
                if rtl:
                    docx_set_rtl_paragraph(p, justify=False)

    for item in sections:
        if item["heading"]:
            docx_heading_bar(
                document,
                item["heading"].lstrip("#").strip(),
                level=1,
                rtl=rtl,
                fill_hex=fill_h,
            )
        body = item.get("body") or ""
        if body.strip():
            # Whole-body collapsed table only?
            maybe = try_parse_pipe_table_blob(body)
            if maybe is not None and body.count("\n") < 2 and body.count("|") > 6:
                docx_add_table(
                    document,
                    list(maybe.get("headers") or []),
                    list(maybe.get("rows") or []),
                    rtl=rtl,
                )
            else:
                docx_write_rich_body(document, body, rtl=rtl)

    # Structured tables from payload (same as PDF path)
    tables = payload.get("tables")
    if isinstance(tables, list):
        for value in tables:
            if not isinstance(value, dict):
                continue
            heading = str(value.get("heading", ""))
            headers = value.get("headers")
            rows = value.get("rows")
            if heading:
                docx_heading_bar(
                    document, heading, level=2, rtl=rtl, fill_hex=fill_h
                )
            if isinstance(headers, list) and isinstance(rows, list) and headers:
                docx_add_table(
                    document,
                    [str(c) for c in headers],
                    [
                        [str(c) for c in row]
                        for row in rows
                        if isinstance(row, list)
                    ],
                    rtl=rtl,
                )

    citations = payload.get("citations")
    if isinstance(citations, list) and citations:
        docx_heading_bar(
            document, chrome["references"], level=1, rtl=rtl, fill_hex=fill_h
        )
        for value in citations:
            p = document.add_paragraph(str(value), style="List Number")
            if rtl:
                docx_set_rtl_paragraph(p, justify=False)
            else:
                docx_force_justify(p)

    # Critical: section + table + theme RTL so Word does not open as LTR shell
    if rtl:
        docx_apply_document_rtl(document)

    document.save(output)


def _generate_xlsx(output: Path, payload: dict[str, Any]) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    workbook = Workbook()
    workbook.remove(workbook.active)
    sheets = payload.get("sheets")
    if not isinstance(sheets, list):
        sheets = []
    for index, value in enumerate(sheets, 1):
        if not isinstance(value, dict):
            continue
        name = str(value.get("name", f"Sheet{index}"))[:31]
        sheet = workbook.create_sheet(name or f"Sheet{index}")
        rows = value.get("rows")
        if not isinstance(rows, list):
            continue
        for row in rows:
            if isinstance(row, list):
                sheet.append([item if item is not None else "" for item in row])
        for cell in sheet[1] if sheet.max_row else ():
            cell.font = Font(bold=True)
    if not workbook.sheetnames:
        workbook.create_sheet("Sheet")
    workbook.save(output)


def _generate_pptx(output: Path, payload: dict[str, Any]) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = str(payload.get("title", "Presentation"))
    for value in payload.get("slides", []) if isinstance(payload.get("slides"), list) else []:
        if not isinstance(value, dict):
            continue
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = str(value.get("heading", ""))
        frame = slide.placeholders[1].text_frame
        bullets = value.get("bullets")
        lines = bullets if isinstance(bullets, list) else str(value.get("body", "")).splitlines()
        for index, line in enumerate(lines):
            if index == 0:
                frame.text = str(line)
            else:
                frame.add_paragraph().text = str(line)
    presentation.save(output)


def _weasy_pdf(output: Path, source: str, assets: Path) -> None:
    from weasyprint import HTML, default_url_fetcher

    try:
        validate_restricted_render_resources(
            source,
            approved_asset_names=frozenset(
                path.name for path in assets.iterdir() if path.is_file()
            ),
        )
    except Exception as exc:
        raise WorkerError(
            "external_resource_denied",
            "External or unapproved local resources are forbidden during rendering",
        ) from exc
    assets_root = assets.resolve()

    def fetch(url: str, *args: Any, **kwargs: Any) -> Any:
        if url.startswith("data:"):
            return default_url_fetcher(url, *args, **kwargs)
        if url.startswith("file:"):
            candidate = Path(url[5:]).resolve()
            if candidate.is_relative_to(assets_root) and candidate.is_file():
                return default_url_fetcher(url, *args, **kwargs)
        raise WorkerError("external_resource_denied", "External rendering fetch was denied")

    HTML(string=source, base_url=assets.as_uri() + "/", url_fetcher=fetch).write_pdf(output)


def _extract_docx_text(source: Path) -> tuple[str, list[str]]:
    """Pull paragraph + table text from a DOCX for lossy PDF conversion."""

    try:
        from docx import Document
    except ImportError as exc:
        raise WorkerError(
            "renderer_unavailable",
            "python-docx is required for DOCX→PDF text fallback (install python-docx)",
        ) from exc

    warnings: list[str] = []
    document = Document(str(source))
    chunks: list[str] = []
    for paragraph in document.paragraphs:
        text = (paragraph.text or "").strip()
        if text:
            chunks.append(text)
    for table in document.tables:
        for row in table.rows:
            cells = [(cell.text or "").strip() for cell in row.cells]
            line = " | ".join(cell for cell in cells if cell)
            if line:
                chunks.append(line)
    if not chunks:
        warnings.append("DOCX had no extractable text; emitted a titled blank PDF")
        return "", warnings
    return "\n\n".join(chunks), warnings


def _office_text_to_pdf(
    output: Path,
    *,
    title: str,
    body: str,
    warnings: list[str],
) -> None:
    """Render plain extracted text to PDF via the reportlab generator."""

    _generate_pdf(
        output,
        {
            "title": title or "Document",
            "sections": [{"heading": "", "body": body or "(empty document)"}],
            "page_numbers": True,
        },
        warnings,
    )
    warnings.append(
        "Converted via text extraction (reportlab-office); layout, images, and "
        "styles are not preserved. Install LibreOffice (soffice) for high-fidelity "
        "Office→PDF conversion."
    )


def _libreoffice(request: dict[str, Any], source: Path, output: Path) -> None:
    from kazma_core.documents.binaries import find_soffice, run_soffice_cli

    if not find_soffice():
        raise WorkerError(
            "renderer_unavailable",
            "Healthy headless LibreOffice is unavailable (soffice not found on PATH "
            "or under Program Files\\LibreOffice)",
        )
    profile = output.parent / "lo-profile"
    profile.mkdir()
    target = output.suffix.lstrip(".")
    env = {
        key: os.environ[key]
        for key in ("PATH", "SYSTEMROOT", "WINDIR", "TEMP", "TMP")
        if key in os.environ
    }
    try:
        result = run_soffice_cli(
            (
                "--headless",
                "--nologo",
                "--nodefault",
                "--nolockcheck",
                "--norestore",
                f"-env:UserInstallation={profile.as_uri()}",
                "--convert-to",
                target,
                "--outdir",
                str(output.parent),
                str(source),
            ),
            timeout=float(request.get("library_timeout_seconds", 120)),
            cwd=output.parent,
            env=env,
        )
    except FileNotFoundError as exc:
        raise WorkerError(
            "renderer_unavailable",
            "Healthy headless LibreOffice is unavailable",
        ) from exc
    except Exception as exc:
        raise WorkerError(
            "conversion_failed",
            f"Headless LibreOffice conversion failed ({type(exc).__name__})",
        ) from exc
    produced = output.parent / f"{source.stem}.{target}"
    if result.returncode or not produced.is_file():
        raise WorkerError("conversion_failed", "Headless LibreOffice conversion failed")
    if produced != output:
        os.replace(produced, output)


def _render(request: dict[str, Any], output: Path) -> tuple[str, str, list[str]]:
    operation = str(request.get("operation", ""))
    payload = request.get("payload")
    if not isinstance(payload, dict):
        raise WorkerError("invalid_request", "Renderer payload must be an object")
    warnings: list[str] = []
    template = payload.get("_template")
    if template not in (None, "default", "report", "compact"):
        raise WorkerError("unsupported_template", "Requested document template is unavailable")
    renderer = str(request.get("renderer", ""))
    source: Path | None = None
    if request.get("source_path") is not None:
        source = Path(str(request["source_path"])).resolve(strict=True)
        if _sha256(source) != request.get("source_sha256"):
            raise WorkerError("document_changed", "Source changed before rendering")
    if operation.startswith("generate:"):
        target = operation.split(":", 1)[1]
        if target == "markdown":
            output.write_text(_markdown(payload), encoding="utf-8")
        elif target == "html":
            output.write_text(_markdown_html(_markdown(payload)), encoding="utf-8")
        elif target == "pdf":
            _generate_pdf(output, payload, warnings)
        elif target == "docx":
            _generate_docx(output, payload)
        elif target == "xlsx":
            _generate_xlsx(output, payload)
        elif target == "pptx":
            _generate_pptx(output, payload)
        else:
            raise WorkerError("unsupported_operation", "Unsupported generation format")
    elif operation == "convert:markdown:html":
        assert source is not None
        text = source.read_text(encoding="utf-8")
        _safe_html(text)
        output.write_text(_markdown_html(text), encoding="utf-8")
    elif operation == "convert:markdown:docx":
        assert source is not None
        _safe_html(source.read_text(encoding="utf-8"))
        _generate_docx(
            output,
            {
                "title": source.stem,
                "sections": [{"heading": "", "body": source.read_text(encoding="utf-8")}],
            },
        )
    elif operation in {"convert:markdown:pdf", "convert:html:pdf"} and renderer != "reportlab-office":
        assert source is not None
        text = source.read_text(encoding="utf-8")
        if operation.startswith("convert:markdown"):
            text = _markdown_html(text)
        _weasy_pdf(output, text, output.parent / "assets")
    elif renderer == "reportlab-office":
        assert source is not None
        if operation == "convert:docx:pdf":
            body, extract_warnings = _extract_docx_text(source)
            warnings.extend(extract_warnings)
            _office_text_to_pdf(output, title=source.stem, body=body, warnings=warnings)
        elif operation in {
            "convert:markdown:pdf",
            "convert:md:pdf",
            "convert:txt:pdf",
            "convert:text:pdf",
        }:
            try:
                body = source.read_text(encoding="utf-8")
            except UnicodeError as exc:
                raise WorkerError(
                    "invalid_document_encoding",
                    "Text conversion source must be valid UTF-8",
                ) from exc
            _office_text_to_pdf(output, title=source.stem, body=body, warnings=warnings)
        else:
            raise WorkerError(
                "unsupported_operation",
                f"reportlab-office cannot handle {operation}",
            )
    elif renderer == "libreoffice":
        assert source is not None
        _libreoffice(request, source, output)
    else:
        raise WorkerError("unsupported_operation", "Unsupported document conversion")
    return renderer, str(request.get("renderer_version", "1")), warnings


def execute(request_path: Path, result_path: Path) -> int:
    try:
        request = _request(request_path)
        request_path.unlink(missing_ok=True)
        _verify_assets(request, result_path.parent)
        _validate_payload_limits(request)
        output = result_path.parent / str(request.get("output_name", ""))
        if output.parent != result_path.parent or not output.name.startswith("output."):
            raise WorkerError("invalid_request", "Invalid worker output path")
        renderer, version, warnings = _render(request, output)
        if request.get("source_path") is not None:
            source = Path(str(request["source_path"])).resolve(strict=True)
            if _sha256(source) != request.get("source_sha256"):
                raise WorkerError("document_changed", "Source changed during rendering")
        _verify_assets(request, result_path.parent)
        if not output.is_file() or output.stat().st_size <= 0:
            raise WorkerError("empty_output", "Renderer produced no document")
        if output.stat().st_size > int(request.get("max_output_bytes", 256 * 1024 * 1024)):
            output.unlink(missing_ok=True)
            raise WorkerError("document_limit_exceeded", "Rendered document exceeds output limit")
        extension = output.suffix.lower().lstrip(".")
        _write(
            result_path,
            {
                "protocol_version": _PROTOCOL_VERSION,
                "ok": True,
                "code": "ok",
                "message": "Document rendered",
                "renderer": renderer,
                "renderer_version": version,
                "source_sha256": request.get("source_sha256"),
                "output_name": output.name,
                "output_extension": extension,
                "output_mime_type": _MIME[extension],
                "output_size": output.stat().st_size,
                "output_sha256": _sha256(output),
                "warnings": warnings,
            },
        )
        return 0
    except WorkerError as exc:
        _write(
            result_path,
            {
                "protocol_version": _PROTOCOL_VERSION,
                "ok": False,
                "code": exc.code,
                "message": exc.safe_message,
            },
        )
        return 0
    except Exception as exc:
        _write(
            result_path,
            {
                "protocol_version": _PROTOCOL_VERSION,
                "ok": False,
                "code": "renderer_worker_failure",
                "message": f"Document rendering failed safely ({type(exc).__name__})",
            },
        )
        return 1


def main() -> int:
    if len(sys.argv) != 3:
        return 2
    return execute(Path(sys.argv[1]), Path(sys.argv[2]))


if __name__ == "__main__":
    raise SystemExit(main())
