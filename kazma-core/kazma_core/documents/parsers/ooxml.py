"""Office Open XML and converter-backed legacy Office parsers."""

from __future__ import annotations

import shutil
import subprocess
import uuid
from dataclasses import replace
from pathlib import Path

from ..errors import DocumentLimitError, DocumentParseError
from ..models import BlockType, DocumentIR
from ..sniff import sniff_document
from .common import IRBuilder, ParseContext

__all__ = ["DocxParser", "LegacyOfficeParser", "PptxParser", "XlsxParser"]


class DocxParser:
    def parse(self, path: Path, context: ParseContext) -> DocumentIR:
        from docx import Document

        document = Document(str(path))
        builder = IRBuilder(path, context)
        builder.count_images(len(document.inline_shapes))
        blocks: list[tuple[BlockType, str, dict[str, object] | None]] = []
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style = (paragraph.style.name or "").lower()
            block_type = BlockType.HEADING if "heading" in style else BlockType.PARAGRAPH
            blocks.append((block_type, text, {"style": paragraph.style.name or ""}))
        for table_index, table in enumerate(document.tables, start=1):
            rows: list[str] = []
            for row_index, row in enumerate(table.rows, start=1):
                if row_index > context.config.max_rows_per_sheet:
                    builder.warnings.append(
                        f"Table {table_index} rows truncated at the configured limit"
                    )
                    break
                cells = [cell.text.strip() for cell in row.cells]
                builder.count_cells(len(cells))
                rows.append(" | ".join(cells))
            blocks.append(
                (
                    BlockType.TABLE,
                    "\n".join(rows),
                    {"table_index": table_index, "rows": len(rows)},
                )
            )
        builder.add_page(blocks, metadata={"kind": "logical_document"})
        return builder.build(metadata={"logical_page_count": 1})


class XlsxParser:
    def parse(self, path: Path, context: ParseContext) -> DocumentIR:
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=True, data_only=True, keep_links=False)
        try:
            if len(workbook.sheetnames) > context.config.max_sheets:
                raise DocumentLimitError("Workbook sheet count exceeds the configured limit")
            builder = IRBuilder(path, context)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                builder.count_images(len(getattr(sheet, "_images", ()) or ()))
                rows: list[str] = []
                row_count = 0
                for row_count, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    if row_count > context.config.max_rows_per_sheet:
                        builder.warnings.append(
                            f"Sheet {sheet_name} rows truncated at the configured limit"
                        )
                        row_count -= 1
                        break
                    values = ["" if value is None else str(value) for value in row]
                    builder.count_cells(len(values))
                    if any(values):
                        rows.append(" | ".join(values))
                builder.add_page(
                    [
                        (
                            BlockType.TABLE,
                            "\n".join(rows),
                            {"rows": row_count, "sheet_name": sheet_name},
                        )
                    ],
                    metadata={
                        "kind": "sheet",
                        "sheet_name": sheet_name,
                        "row_count": row_count,
                    },
                )
            return builder.build(metadata={"sheet_count": len(workbook.sheetnames)})
        finally:
            workbook.close()


class PptxParser:
    def parse(self, path: Path, context: ParseContext) -> DocumentIR:
        from pptx import Presentation

        presentation = Presentation(str(path))
        if len(presentation.slides) > context.config.max_slides:
            raise DocumentLimitError("Presentation slide count exceeds the configured limit")
        builder = IRBuilder(path, context)
        for slide_index, slide in enumerate(presentation.slides, start=1):
            blocks: list[tuple[BlockType, str, dict[str, object] | None]] = []
            for shape_index, shape in enumerate(slide.shapes, start=1):
                if int(getattr(shape, "shape_type", 0) or 0) == 13:
                    builder.count_images(1)
                if getattr(shape, "has_text_frame", False):
                    text = "\n".join(
                        paragraph.text.strip()
                        for paragraph in shape.text_frame.paragraphs
                        if paragraph.text.strip()
                    )
                    if text:
                        blocks.append(
                            (BlockType.TEXT, text, {"shape_index": shape_index})
                        )
                if getattr(shape, "has_table", False):
                    rows = []
                    for row_index, row in enumerate(shape.table.rows, start=1):
                        if row_index > context.config.max_rows_per_sheet:
                            builder.warnings.append(
                                f"Slide {slide_index} table rows truncated"
                            )
                            break
                        cells = [cell.text.strip() for cell in row.cells]
                        builder.count_cells(len(cells))
                        rows.append(" | ".join(cells))
                    blocks.append(
                        (
                            BlockType.TABLE,
                            "\n".join(rows),
                            {"shape_index": shape_index, "rows": len(rows)},
                        )
                    )
            builder.add_page(
                blocks,
                metadata={"kind": "slide", "slide_number": slide_index},
            )
        return builder.build(metadata={"slide_count": len(presentation.slides)})


class LegacyOfficeParser:
    """Convert legacy OLE Office files with a probed headless LibreOffice."""

    def __init__(self, executable: str) -> None:
        self.executable = executable

    def parse(self, path: Path, context: ParseContext) -> DocumentIR:
        output_extension = {".doc": ".docx", ".xls": ".xlsx", ".ppt": ".pptx"}[
            path.suffix.lower()
        ]
        run_root = context.config.storage_root / "parser-runs"
        run_dir = run_root / f"convert-{uuid.uuid4().hex}"
        run_dir.mkdir(parents=True, exist_ok=False)
        try:
            process = subprocess.run(
                [
                    self.executable,
                    "--headless",
                    "--convert-to",
                    output_extension.lstrip("."),
                    "--outdir",
                    str(run_dir),
                    str(path),
                ],
                capture_output=True,
                check=False,
                timeout=max(5, context.config.worker_timeout_seconds - 5),
                text=True,
            )
            converted = run_dir / f"{path.stem}{output_extension}"
            if process.returncode != 0 or not converted.is_file():
                raise DocumentParseError(
                    "LibreOffice could not safely convert the legacy document",
                    code="legacy_conversion_failed",
                )
            sniffed = sniff_document(converted, context.config)
            converted_context = replace(
                context,
                mime_type=sniffed.mime_type,
                extension=sniffed.extension,
            )
            parser = {
                ".docx": DocxParser,
                ".xlsx": XlsxParser,
                ".pptx": PptxParser,
            }[output_extension]()
            parsed = parser.parse(converted, converted_context)
            metadata = dict(parsed.metadata)
            metadata["legacy_source_extension"] = path.suffix.lower()
            metadata["source_name"] = path.name
            return replace(parsed, metadata=metadata)
        finally:
            shutil.rmtree(run_dir, ignore_errors=True)
