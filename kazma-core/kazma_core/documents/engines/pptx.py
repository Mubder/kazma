"""PPTX engine for the unified document layer — branded layout.

Slides don't share the document :class:`ContentModel` (they have their own
slides/bullets shape), so this engine consumes the native ``slides`` payload
**plus** a :class:`~kazma_core.documents.profile.DocProfile`. The profile
supplies the shared theme (accent / heading colours, fonts) and direction.

The deck is built on a blank layout with explicit shapes so the look is fully
controlled and consistent with the DOCX/PDF/HTML family:

  - 16:9 widescreen (modern default).
  - Title slide: full-width accent band with the title (white, centred).
  - Content slides: a **heading bar** (accent fill + white title) — the same
    heading-bar motif as DOCX/PDF/HTML — followed by the bullet body and a
    footer (brand + page number).
  - RTL: every paragraph carries ``rtl="1"`` AND ``algn="r"`` (rtl alone does
    not right-align; the DrawingML attribute is ``algn``, not ``al``).
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from kazma_core.documents.profile import DocProfile

logger = logging.getLogger(__name__)

__all__ = ["PptxEngine"]


class PptxEngine:
    """Render a slides payload to a branded ``.pptx`` under a :class:`DocProfile`."""

    def __init__(self, profile: DocProfile) -> None:
        self.profile = profile
        self.theme = profile.theme

    # ------------------------------------------------------------------ #
    # public entry
    # ------------------------------------------------------------------ #
    def render(self, payload: dict[str, Any], output: Path | str) -> None:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        t = self.theme
        colors = {
            "accent": str(t["accent"]).lstrip("#"),
            "heading": str(t["heading_fill"]).lstrip("#"),
            "body": str(t["body"]).lstrip("#"),
            "muted": str(t["muted"]).lstrip("#"),
        }

        presentation = Presentation()
        # 16:9 widescreen.
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        # python-pptx leaves a stale type="screen4x3" label after resizing;
        # correct it so viewers categorise the deck as 16:9 (cosmetic but right).
        try:
            presentation._sldSz.set("type", "screen16x9")
        except Exception:
            logger.debug("[pptx] could not set sldSz type", exc_info=True)
        slide_w = presentation.slide_width
        slide_h = presentation.slide_height
        blank = presentation.slide_layouts[6]

        # File core properties (title/author/subject).
        try:
            cp = presentation.core_properties
            cp.title = str(payload.get("title") or "Presentation")
            cp.author = str(payload.get("author") or "Kazma")
            if payload.get("subject"):
                cp.subject = str(payload.get("subject"))
            if payload.get("keywords"):
                cp.keywords = str(payload.get("keywords"))
        except Exception:
            logger.debug("[pptx] core properties failed", exc_info=True)

        # ── Title slide ─────────────────────────────────────────────────── #
        title_slide = presentation.slides.add_slide(blank)
        band = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, slide_w, Inches(2.8),
        )
        self._fill(band, colors["accent"])
        self._shape_text(
            band, str(payload.get("title", "Presentation")),
            color_hex="FFFFFF", bold=True, size=36, align="ctr", band=True,
        )
        subtitle = payload.get("subtitle")
        if subtitle:
            sub = title_slide.shapes.add_textbox(
                Inches(0.8), Inches(3.1), slide_w - Inches(1.6), Inches(0.8),
            )
            self._shape_text(sub, str(subtitle), color_hex=colors["muted"],
                             size=20, align="ctr")
        # Optional speaker notes on the title slide.
        self._set_notes(title_slide, payload.get("notes"))

        # ── Content slides ─────────────────────────────────────────────── #
        slides = payload.get("slides") if isinstance(payload.get("slides"), list) else []
        page_num = 1
        for value in slides:
            if not isinstance(value, dict):
                continue
            page_num += 1
            slide = presentation.slides.add_slide(blank)

            # Heading bar (mirrors DOCX/PDF/HTML heading bars).
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, Inches(0.25), slide_w, Inches(0.95),
            )
            self._fill(bar, colors["heading"])
            self._shape_text(
                bar, str(value.get("heading", "")),
                color_hex="FFFFFF", bold=True, size=26, align=None, band=True,
            )

            # Bullet body.
            body = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.55),
                slide_w - Inches(1.2), slide_h - Inches(2.5),
            )
            body.fill.background()  # transparent
            body.line.fill.background()
            self._bullet_body(body, value, colors["body"])

            # Footer (brand + page number).
            footer = slide.shapes.add_textbox(
                Inches(0.4), slide_h - Inches(0.45), slide_w - Inches(0.8), Inches(0.3),
            )
            chrome = self.profile.chrome
            page_fmt = str(chrome.get("page_fmt") or "{n}")
            self._shape_text(
                footer,
                f'{chrome["brand"]}    ·    {page_fmt.format(n=page_num)}',
                color_hex=colors["muted"], size=10, align=None,
            )
            # Optional speaker notes for this slide.
            self._set_notes(slide, value.get("notes"))

        presentation.save(str(output))

    @staticmethod
    def _set_notes(slide: Any, notes: Any) -> None:
        """Attach speaker notes to a slide (creates the notes slide if needed).

        ``notes`` may be a string or a list of bullet strings.
        """
        if not notes:
            return
        lines = notes if isinstance(notes, list) else [str(notes)]
        tf = slide.notes_slide.notes_text_frame
        tf.text = lines[0]
        for line in lines[1:]:
            tf.add_paragraph().text = str(line)

    # ================================================================== #
    # shape helpers
    # ================================================================== #
    def _fill(self, shape: Any, hex_color: str) -> None:
        from pptx.dml.color import RGBColor

        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(hex_color)
        shape.line.fill.background()  # no border

    def _shape_text(self, shape: Any, text: str, *, color_hex: str, bold: bool = False,
                    size: float = 18, align: str | None = None, band: bool = False) -> None:
        """Put themed text in a shape and style its single paragraph.

        ``align`` defaults to right (RTL) / left (LTR); ``band`` tightens the
        internal margins and vertically centres text for full-width bars.
        """
        from pptx.enum.text import MSO_ANCHOR
        from pptx.util import Inches, Pt

        tf = shape.text_frame
        tf.word_wrap = True
        if band:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            for side in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
                setattr(tf, side, Inches(0.15))
        tf.text = text
        self._style_paragraph(tf.paragraphs[0], color_hex=color_hex, bold=bold,
                              size=size, align=align)

    def _bullet_body(self, shape: Any, value: dict[str, Any], body_hex: str) -> None:
        """Fill a body shape with bullets (or body lines), each a themed RTL/LTR paragraph."""
        tf = shape.text_frame
        tf.word_wrap = True
        bullets = value.get("bullets")
        lines = bullets if isinstance(bullets, list) else str(value.get("body", "")).splitlines()
        for index, line in enumerate(lines):
            para = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            self._style_paragraph(para, text=f"•  {line}", color_hex=body_hex, size=18)

    def _style_paragraph(self, paragraph: Any, *, text: str | None = None, color_hex: str,
                         bold: bool = False, size: float = 18,
                         align: str | None = None) -> None:
        """Set a paragraph's optional text + themed run + RTL direction + alignment.

        RTL needs ``rtl="1"`` (reading order) AND ``algn`` alignment. ``rtl`` alone
        does not right-align — the paragraph inherits left from the layout, so
        Arabic looks LTR. DrawingML's attribute is ``algn`` (NOT ``al``).
        """
        from pptx.dml.color import RGBColor
        from pptx.util import Pt

        if text is not None:
            paragraph.text = text
        pPr = paragraph._p.get_or_add_pPr()
        pPr.set("rtl", "1" if self.profile.rtl else "0")
        if align is None:
            align = "r" if self.profile.rtl else "l"
        pPr.set("algn", align)
        if paragraph.runs:
            run = paragraph.runs[0]
            run.font.size = Pt(size)
            run.font.name = "Calibri"
            run.font.bold = bold
            run.font.color.rgb = RGBColor.from_string(color_hex)
