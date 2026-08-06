"""Attachment → LLM message-content builder.

Shared by the gateway path (``agent_handler/store.py``) and the Web SSE
path (``kazma_ui/sse_chat.py``) so both transports produce identical
OpenAI-compatible multimodal content from an :class:`Attachment` list.

Policy (see roadmap Phase 1.2):

* **Images** (PNG/JPEG/WEBP/GIF, ≤ ``MAX_INLINE_BYTES``) are inlined as
  base64 ``data:`` URIs into the ``image_url`` content block so the LLM
  sees them immediately. This mirrors the proven pattern in
  ``kazma_core/tools/vision_analyze.py:_build_vision_messages``.
* **Documents / large media / over the inline cap** are persisted to
  ``kazma-data/attachments/`` and represented in the prompt as a text
  stub pointing the agent at the file via ``file_read``. This keeps the
  prompt size bounded while still making the bytes reachable.
"""

from __future__ import annotations

import base64
import logging
import uuid
from pathlib import Path
from typing import TYPE_CHECKING, Any

if TYPE_CHECKING:
    from kazma_gateway.gateway import Attachment

logger = logging.getLogger(__name__)

# Inlining images larger than this as base64 bloats the prompt. The vision
# tool uses a 20 MB cap; we are more conservative for chat to avoid context
# blow-up. Larger images are persisted and referenced as files.
MAX_INLINE_BYTES = 8 * 1024 * 1024  # 8 MB

# MIME types safe to inline as vision input to OpenAI-compatible providers.
# Matches vision_analyze.py's accepted set.
_INLINE_IMAGE_MIMES = frozenset(
    {"image/png", "image/jpeg", "image/jpg", "image/webp", "image/gif"}
)

# Where over-cap / non-image attachments are persisted so the agent can
# open them with file_read. Relative to CWD, matching tools/image_gen.py.
ATTACHMENT_DIR = Path("kazma-data/attachments")


def _persist_attachment(attachment: "Attachment") -> str:
    """Persist attachment bytes to disk and return the absolute file path.

    Assumes ``attachment.data`` is populated. The filename is uniqueified
    to avoid collisions across concurrent turns.
    """
    ATTACHMENT_DIR.mkdir(parents=True, exist_ok=True)
    base = attachment.filename or f"{attachment.kind}_{uuid.uuid4().hex[:8]}"
    # Disambiguate while preserving extension.
    dest = ATTACHMENT_DIR / f"{Path(base).stem}_{uuid.uuid4().hex[:6]}{Path(base).suffix}"
    dest.write_bytes(attachment.data or b"")
    return str(dest.resolve())


def _try_parse_attachment(path: str, *, max_chars: int = 2000) -> str:
    """Best-effort synchronous parse of a document attachment for prompt injection.

    Called by build_user_content when a document (PDF/DOCX/XLSX/PPTX/CSV) is
    uploaded — extracts the first ``max_chars`` of text so the agent sees the
    content immediately in the prompt. Returns an error string on failure.

    This is synchronous (not async) because build_user_content is sync and
    runs in the graph-builder thread before the LLM call.
    """
    try:
        p = Path(path)
        suffix = p.suffix.lower()

        if suffix == ".pdf":
            try:
                import pdfplumber

                with pdfplumber.open(p) as pdf:
                    lines = []
                    for i, page in enumerate(pdf.pages[:5]):
                        text = page.extract_text() or ""
                        if text.strip():
                            lines.append(text.strip())
                        if sum(len(l) for l in lines) > max_chars:
                            break
                    result = "\n".join(lines)[:max_chars]
                    return result if result.strip() else "Error: no text layer (scanned PDF)"
            except ImportError:
                return "Error: pdfplumber not installed"

        elif suffix == ".docx":
            try:
                from docx import Document

                doc = Document(str(p))
                lines = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
                return "\n".join(lines)[:max_chars]
            except ImportError:
                return "Error: python-docx not installed"

        elif suffix in (".xlsx", ".xls"):
            try:
                from openpyxl import load_workbook

                wb = load_workbook(p, read_only=True, data_only=True)
                lines = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    lines.append(f"--- {sheet_name} ---")
                    for i, row in enumerate(ws.iter_rows(values_only=True)):
                        if i >= 20:
                            break
                        cells = [str(c) if c else "" for c in row]
                        if any(c.strip() for c in cells):
                            lines.append(" | ".join(cells))
                wb.close()
                return "\n".join(lines)[:max_chars]
            except ImportError:
                return "Error: openpyxl not installed"

        elif suffix == ".csv":
            import csv as _csv

            with open(p, "r", encoding="utf-8", errors="replace", newline="") as f:
                reader = _csv.reader(f)
                lines = [", ".join(row) for _, row in zip(range(20), reader)]
            return "\n".join(lines)[:max_chars]

        elif suffix == ".pptx":
            try:
                from pptx import Presentation

                prs = Presentation(str(p))
                lines = []
                for i, slide in enumerate(prs.slides[:10]):
                    lines.append(f"--- Slide {i+1} ---")
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                if para.text.strip():
                                    lines.append(para.text.strip())
                return "\n".join(lines)[:max_chars]
            except ImportError:
                return "Error: python-pptx not installed"

        else:
            return f"Error: unsupported format {suffix}"

    except Exception as exc:
        logger.debug("[attachments] document parse failed: %s", exc)
        return f"Error: {exc}"


def _resolve_active_model_vision_capable() -> bool:
    """Best-effort: is the active model vision-capable?

    Returns ``True`` (fail-open, inline images as before) on any error so
    that an inability to read the registry never silently drops images.
    """
    try:
        from kazma_core.model_registry import get_model_registry
        from kazma_core.vision_capability import active_model_is_vision_capable

        return active_model_is_vision_capable(get_model_registry())
    except Exception as exc:  # noqa: BLE001 — best-effort
        logger.debug(
            "[attachments] could not resolve active model vision capability: %s",
            exc,
        )
        return True  # fail-open


def build_user_content(
    text: str,
    attachments: list["Attachment"] | None,
    vision_capable: bool | None = None,
) -> str | list[dict[str, Any]]:
    """Build the OpenAI ``content`` for a user turn.

    Returns a plain string when there are no inlinable attachments (the
    fast path every plain-text message takes), or a multimodal
    ``content`` list (text + image_url blocks) when images are present.

    Non-inlinable attachments (documents, audio, over-cap images) are
    persisted and folded into the text portion as ``[Attached: <path>]``
    stubs so the agent can fetch them with ``file_read``.

    ``vision_capable`` controls whether inlinable images are actually
    inlined as ``image_url`` blocks. When the active model is text-only
    (e.g. DeepSeek) passing ``vision_capable=False`` forces images down
    the persist-and-stub path so the provider never receives an
    ``image_url`` part it would reject with a 400. When ``None`` (default),
    the active model is looked up best-effort via
    :func:`kazma_core.vision_capability.active_model_is_vision_capable`;
    on any error, behavior defaults to inlining (fail-open).
    """
    if not attachments:
        return text

    # Resolve the active model's vision capability if the caller didn't.
    if vision_capable is None:
        vision_capable = _resolve_active_model_vision_capable()

    blocks: list[dict[str, Any]] = []
    text_parts: list[str] = [text] if text else []
    saw_image = False

    for att in attachments:
        data = att.data
        # Fetch on demand if only a URL was provided.
        if data is None and att.url:
            try:
                import httpx

                resp = httpx.get(att.url, timeout=30.0, follow_redirects=True)
                resp.raise_for_status()
                data = resp.content
            except Exception as exc:  # noqa: BLE001 — network is best-effort
                logger.warning(
                    "[attachments] failed to fetch %s: %s", att.url, exc
                )
                text_parts.append(
                    f"[Attached: {att.filename or att.url} — fetch failed]"
                )
                continue

        is_inlinable_image = (
            att.kind == "image"
            and att.mime in _INLINE_IMAGE_MIMES
            and data is not None
            and len(data) <= MAX_INLINE_BYTES
        )
        # A text-only model cannot consume an image_url content part —
        # force the image down the persist-and-stub path (file_read) so
        # the provider receives plain text and never 400s on `image_url`.
        if is_inlinable_image and not vision_capable:
            is_inlinable_image = False
            logger.info(
                "[attachments] active model is text-only — downgrading "
                "image %s to file reference (no image_url)",
                att.filename or att.url,
            )

        if is_inlinable_image and data is not None:
            data_uri = (
                f"data:{att.mime};base64,{base64.b64encode(data).decode('ascii')}"
            )
            blocks.append(
                {"type": "image_url", "image_url": {"url": data_uri}}
            )
            saw_image = True
        else:
            # Persist and reference. Non-image or oversized attachments
            # stay out of the prompt payload themselves.
            if data is not None:
                try:
                    path = _persist_attachment(att)
                    # ── Auto-parse document attachments ──────────────
                    # For document types (PDF, DOCX, XLSX, PPTX, CSV), extract
                    # text immediately so the agent sees the content in the
                    # prompt — not just a "use file_read" stub that would
                    # fail on binary formats. The file is also persisted so
                    # the agent can use read_document for the full content.
                    from pathlib import Path as _Path

                    _suffix = _Path(path).suffix.lower()
                    _DOC_SUFFIXES = frozenset({
                        ".pdf", ".docx", ".doc", ".xlsx", ".xls",
                        ".pptx", ".ppt", ".csv", ".json", ".rtf",
                    })
                    if _suffix in _DOC_SUFFIXES:
                        excerpt = _try_parse_attachment(path, max_chars=2000)
                        if excerpt and not excerpt.startswith("Error"):
                            text_parts.append(
                                f"[Attached: {att.filename or path} ({att.mime}) "
                                f"— parsed contents:\n{excerpt}]"
                            )
                            text_parts.append(
                                f"[Full document at: {path} — use read_document for complete content]"
                            )
                        else:
                            text_parts.append(
                                f"[Attached: {att.filename or path} ({att.mime}) "
                                f"— use read_document to open: {path}]"
                            )
                    else:
                        text_parts.append(
                            f"[Attached: {att.filename or path} ({att.mime}) "
                            f"— use file_read to open: {path}]"
                        )
                except Exception as exc:  # noqa: BLE001
                    logger.warning("[attachments] persist failed: %s", exc)
                    text_parts.append(
                        f"[Attached: {att.filename} — save failed: {exc}]"
                    )
            else:
                # No data and no fetchable URL — record what we know.
                label = att.filename or att.url or att.kind
                text_parts.append(f"[Attached: {label} ({att.mime}) — unavailable]")

    if not saw_image:
        # No vision input — keep content a plain string for token efficiency
        # and broad provider compatibility.
        return "\n".join(text_parts)

    # Multimodal: text block first (the user's caption/intent), then images.
    combined_text = "\n".join(text_parts)
    content: list[dict[str, Any]] = []
    if combined_text.strip():
        content.append({"type": "text", "text": combined_text})
    content.extend(blocks)
    return content
